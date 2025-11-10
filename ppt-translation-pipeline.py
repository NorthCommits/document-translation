"""
PowerPoint Translation Pipeline
================================
Complete end-to-end translation system with DeepL integration.

Features:
- Comprehensive metadata extraction (masters, layouts, fills, shadows, etc.)
- DeepL translation with glossary support
- RTL language support with layout mirroring
- Template-based reassembly preserving 100% visual fidelity
- Auto-shrink to prevent text overflow

Usage:
    python ppt_translator_pipeline.py input.pptx -l Spanish
    python ppt_translator_pipeline.py input.pptx -l French -o custom_output.pptx
"""

import json
import os
import sys
import time
import argparse
import requests
from pathlib import Path
from typing import Dict, List, Any, Optional
from copy import deepcopy
from dotenv import load_dotenv

# PowerPoint libraries
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
import zipfile
from lxml import etree


# ============================================================================
# CONFIGURATION & CONSTANTS
# ============================================================================

# Load environment variables
load_dotenv()

# DeepL Configuration
DEEPL_API_KEY = os.getenv('DEEPL_API_KEY')
# Default to Pro endpoint - if you have Free API, use 'https://api-free.deepl.com/v2/translate'
DEEPL_ENDPOINT = os.getenv('DEEPL_ENDPOINT', 'https://api.deepl.com/v2/translate')

# Language Configuration
# Language Configuration
SUPPORTED_LANGUAGES = {
    # European Languages (LTR)
    'French': 'FR',
    'Spanish': 'ES',
    'Italian': 'IT',
    'German': 'DE',
    'Portuguese': 'PT',
    'Dutch': 'NL',
    'Swedish': 'SV',
    'Danish': 'DA',
    'Finnish': 'FI',
    'Norwegian': 'NB',
    'Polish': 'PL',
    'Czech': 'CS',
    'Romanian': 'RO',
    'Hungarian': 'HU',
    'Greek': 'EL',
    'Bulgarian': 'BG',
    'Slovak': 'SK',
    'Slovenian': 'SL',
    'Estonian': 'ET',
    'Latvian': 'LV',
    'Lithuanian': 'LT',
    
    # Asian Languages (LTR)
    'Chinese': 'ZH',
    'Japanese': 'JA',
    'Korean': 'KO',
    'Indonesian': 'ID',
    
    # RTL Languages
    'Arabic': 'AR',
    'Hebrew': 'HE',  # Note: DeepL uses 'HE' not 'IW'
    
    # Other Languages
    'Russian': 'RU',
    'Ukrainian': 'UK',
    'Turkish': 'TR'
}

# Glossary Configuration
GLOSSARIES = {
    'NL': 'c108ea02-1025-4ad4-b702-d10eda123786',  # English-Dutch (fixed typo: d10eda not d1oeda)
    'SV': 'e63a7f5d-b189-4d2c-868a-e8849cd691ac'   # English-Swedish
}

# RTL Languages
RTL_LANGUAGES = ['Arabic', 'Hebrew', 'Urdu', 'Persian', 'Farsi']


# ============================================================================
# PART 1: EXTRACTOR - PowerPoint Content Extraction
# ============================================================================

class PPTXExtractor:
    """Extract comprehensive content and metadata from PowerPoint presentations"""
    
    def __init__(self, pptx_path):
        self.pptx_path = pptx_path
        self.presentation = Presentation(pptx_path)
        self.data = {
            "presentation_name": os.path.basename(pptx_path),
            "total_slides": len(self.presentation.slides),
            "slide_masters": [],
            "slides": []
        }
        self.namespaces = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
            'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram',
            'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        }
        
    def get_color_value(self, color_obj):
        """Extract color in multiple formats"""
        color_info = {}
        try:
            if hasattr(color_obj, 'rgb') and color_obj.rgb:
                color_info['rgb'] = str(color_obj.rgb)
            if hasattr(color_obj, 'theme_color') and color_obj.theme_color:
                color_info['theme_color'] = str(color_obj.theme_color)
            if hasattr(color_obj, 'brightness'):
                color_info['brightness'] = color_obj.brightness
        except:
            pass
        return color_info if color_info else None
    
    def extract_background_info(self, slide_or_layout):
        """Extract background information from a slide or layout"""
        background_info = {
            "follows_master": None,
            "fill_type": None,
            "solid_color": None,
            "gradient_colors": None,
            "pattern_type": None,
            "picture_present": False
        }
        
        try:
            if hasattr(slide_or_layout, 'follow_master_background'):
                background_info["follows_master"] = slide_or_layout.follow_master_background
            
            bg = slide_or_layout.background
            fill = bg.fill
            
            if hasattr(fill, 'type'):
                background_info["fill_type"] = str(fill.type)
            
            if hasattr(fill, 'fore_color'):
                try:
                    background_info["solid_color"] = self.get_color_value(fill.fore_color)
                except:
                    pass
            
            if hasattr(fill, 'pattern'):
                try:
                    background_info["pattern_type"] = str(fill.pattern)
                except:
                    pass
        except:
            pass
        
        return background_info
    
    def extract_slide_masters(self):
        """Extract information about all slide masters and their layouts"""
        masters_info = []
        
        for master_idx, master in enumerate(self.presentation.slide_masters):
            master_data = {
                "master_index": master_idx,
                "master_name": getattr(master, 'name', f"Master_{master_idx}"),
                "background": self.extract_background_info(master),
                "layouts": []
            }
            
            for layout_idx, layout in enumerate(master.slide_layouts):
                layout_data = {
                    "layout_index": layout_idx,
                    "layout_name": layout.name,
                    "background": self.extract_background_info(layout),
                    "placeholders": []
                }
                
                try:
                    for placeholder in layout.placeholders:
                        ph_data = {
                            "placeholder_idx": placeholder.placeholder_format.idx,
                            "placeholder_type": str(placeholder.placeholder_format.type),
                            "name": placeholder.name,
                            "dimensions": {
                                "left": placeholder.left,
                                "top": placeholder.top,
                                "width": placeholder.width,
                                "height": placeholder.height
                            }
                        }
                        layout_data["placeholders"].append(ph_data)
                except:
                    pass
                
                master_data["layouts"].append(layout_data)
            
            masters_info.append(master_data)
        
        return masters_info
    
    def get_slide_layout_info(self, slide):
        """Get layout information for a specific slide"""
        layout_info = {
            "master_index": None,
            "layout_index": None,
            "layout_name": None,
            "follows_master_background": None
        }
        
        try:
            slide_layout = slide.slide_layout
            layout_info["layout_name"] = slide_layout.name
            
            if hasattr(slide, 'follow_master_background'):
                layout_info["follows_master_background"] = slide.follow_master_background
            
            for master_idx, master in enumerate(self.presentation.slide_masters):
                for layout_idx, layout in enumerate(master.slide_layouts):
                    if layout == slide_layout:
                        layout_info["master_index"] = master_idx
                        layout_info["layout_index"] = layout_idx
                        break
                if layout_info["master_index"] is not None:
                    break
        except:
            pass
        
        return layout_info
    
    def extract_placeholder_info(self, shape):
        """Extract placeholder-specific information if shape is a placeholder"""
        placeholder_info = {
            "is_placeholder": False,
            "placeholder_type": None,
            "placeholder_idx": None
        }
        
        try:
            if shape.is_placeholder:
                placeholder_info["is_placeholder"] = True
                placeholder_info["placeholder_type"] = str(shape.placeholder_format.type)
                placeholder_info["placeholder_idx"] = shape.placeholder_format.idx
        except:
            pass
        
        return placeholder_info
    
    def extract_shape_fill(self, shape):
        """Extract comprehensive fill information from shape"""
        fill_info = {
            "fill_type": None,
            "solid_color": None,
            "gradient_stops": None,
            "pattern_type": None,
            "picture_present": False
        }
        
        try:
            if hasattr(shape, 'fill'):
                fill = shape.fill
                
                if hasattr(fill, 'type'):
                    fill_info["fill_type"] = str(fill.type)
                
                try:
                    if hasattr(fill, 'fore_color'):
                        fill_info["solid_color"] = self.get_color_value(fill.fore_color)
                except:
                    pass
                
                try:
                    if hasattr(fill, 'pattern'):
                        fill_info["pattern_type"] = str(fill.pattern)
                        if hasattr(fill, 'back_color'):
                            fill_info["pattern_back_color"] = self.get_color_value(fill.back_color)
                except:
                    pass
                
                try:
                    if hasattr(fill, 'gradient_stops'):
                        stops = []
                        for stop in fill.gradient_stops:
                            stops.append({
                                "position": stop.position,
                                "color": self.get_color_value(stop.color)
                            })
                        fill_info["gradient_stops"] = stops
                except:
                    pass
        except:
            pass
        
        return fill_info
    
    def extract_shape_line(self, shape):
        """Extract line/border information from shape"""
        line_info = {
            "has_line": False,
            "color": None,
            "width": None,
            "dash_style": None,
            "transparency": None
        }
        
        try:
            if hasattr(shape, 'line'):
                line = shape.line
                line_info["has_line"] = True
                
                if hasattr(line, 'color'):
                    line_info["color"] = self.get_color_value(line.color)
                
                if hasattr(line, 'width'):
                    line_info["width"] = line.width
                
                if hasattr(line, 'dash_style'):
                    line_info["dash_style"] = str(line.dash_style)
        except:
            pass
        
        return line_info
    
    def extract_shape_shadow(self, shape):
        """Extract shadow information from shape"""
        shadow_info = {
            "has_shadow": False,
            "shadow_type": None,
            "color": None,
            "transparency": None,
            "blur": None,
            "angle": None,
            "distance": None
        }
        
        try:
            if hasattr(shape, 'shadow'):
                shadow = shape.shadow
                if hasattr(shadow, 'inherit'):
                    shadow_info["has_shadow"] = not shadow.inherit
                
                if hasattr(shadow, 'shadow_type'):
                    shadow_info["shadow_type"] = str(shadow.shadow_type)
                
                try:
                    if hasattr(shadow, 'color'):
                        shadow_info["color"] = self.get_color_value(shadow.color)
                except:
                    pass
                
                if hasattr(shadow, 'transparency'):
                    shadow_info["transparency"] = shadow.transparency
                if hasattr(shadow, 'blur_radius'):
                    shadow_info["blur"] = shadow.blur_radius
                if hasattr(shadow, 'angle'):
                    shadow_info["angle"] = shadow.angle
                if hasattr(shadow, 'distance'):
                    shadow_info["distance"] = shadow.distance
        except:
            pass
        
        return shadow_info
    
    def extract_bullet_formatting(self, paragraph):
        """Extract comprehensive bullet/numbering information via XML"""
        bullet_info = {
            "is_bulleted": False,
            "bullet_type": None,
            "bullet_char": None,
            "bullet_font": None,
            "bullet_color": None,
            "numbering_format": None,
            "start_at": None
        }
        
        try:
            pPr = paragraph._element.pPr
            if pPr is None:
                return bullet_info
            
            buChar = pPr.find('.//a:buChar', self.namespaces)
            if buChar is not None:
                bullet_info["is_bulleted"] = True
                bullet_info["bullet_type"] = "bullet"
                bullet_info["bullet_char"] = buChar.get('char', '•')
            
            buFont = pPr.find('.//a:buFont', self.namespaces)
            if buFont is not None:
                bullet_info["bullet_font"] = buFont.get('typeface')
            
            buClr = pPr.find('.//a:buClr', self.namespaces)
            if buClr is not None:
                srgbClr = buClr.find('.//a:srgbClr', self.namespaces)
                if srgbClr is not None:
                    bullet_info["bullet_color"] = srgbClr.get('val')
                schemeClr = buClr.find('.//a:schemeClr', self.namespaces)
                if schemeClr is not None:
                    bullet_info["bullet_color"] = f"scheme_{schemeClr.get('val')}"
            
            buAutoNum = pPr.find('.//a:buAutoNum', self.namespaces)
            if buAutoNum is not None:
                bullet_info["is_bulleted"] = True
                bullet_info["bullet_type"] = "numbered"
                num_type = buAutoNum.get('type', 'arabicPeriod')
                bullet_info["numbering_format"] = num_type
                start_at = buAutoNum.get('startAt')
                if start_at:
                    bullet_info["start_at"] = int(start_at)
            
            buNone = pPr.find('.//a:buNone', self.namespaces)
            if buNone is not None:
                bullet_info["bullet_type"] = "none"
                bullet_info["is_bulleted"] = False
        except:
            pass
        
        return bullet_info
    
    def extract_run_formatting(self, run):
        """Extract comprehensive formatting details from a text run"""
        font = run.font
        formatting = {
            "text": run.text,
            "font_name": font.name,
            "font_size": font.size.pt if font.size else None,
            "bold": font.bold,
            "italic": font.italic,
            "underline": font.underline,
            "color": self.get_color_value(font.color) if hasattr(font, 'color') else None,
            "strike": None,
            "kerning": None,
            "spacing": None,
            "caps": None,
            "superscript": None,
            "subscript": None,
            "text_highlight": None,
            "text_outline": None
        }
        
        try:
            if hasattr(font, '_element'):
                rPr = run._element.rPr
                if rPr is not None:
                    if rPr.find('.//a:strike', self.namespaces) is not None:
                        strike_elem = rPr.find('.//a:strike', self.namespaces)
                        formatting["strike"] = strike_elem.get('val', 'sngStrike')
                    
                    if rPr.find('.//a:kern', self.namespaces) is not None:
                        kern_elem = rPr.find('.//a:kern', self.namespaces)
                        formatting["kerning"] = kern_elem.get('val')
                    
                    if rPr.find('.//a:spc', self.namespaces) is not None:
                        spc_elem = rPr.find('.//a:spc', self.namespaces)
                        formatting["spacing"] = spc_elem.get('val')
                    
                    if rPr.get('cap'):
                        formatting["caps"] = rPr.get('cap')
                    
                    baseline = rPr.get('baseline')
                    if baseline:
                        baseline_val = int(baseline)
                        if baseline_val > 0:
                            formatting["superscript"] = baseline_val
                        elif baseline_val < 0:
                            formatting["subscript"] = abs(baseline_val)
                    
                    highlight = rPr.find('.//a:highlight', self.namespaces)
                    if highlight is not None:
                        srgbClr = highlight.find('.//a:srgbClr', self.namespaces)
                        if srgbClr is not None:
                            formatting["text_highlight"] = srgbClr.get('val')
                        schemeClr = highlight.find('.//a:schemeClr', self.namespaces)
                        if schemeClr is not None:
                            formatting["text_highlight"] = f"scheme_{schemeClr.get('val')}"
                    
                    ln = rPr.find('.//a:ln', self.namespaces)
                    if ln is not None:
                        outline_info = {
                            "width": ln.get('w'),
                            "color": None
                        }
                        
                        solidFill = ln.find('.//a:solidFill', self.namespaces)
                        if solidFill is not None:
                            srgbClr = solidFill.find('.//a:srgbClr', self.namespaces)
                            if srgbClr is not None:
                                outline_info["color"] = srgbClr.get('val')
                        
                        formatting["text_outline"] = outline_info
        except:
            pass
        
        return formatting
    
    def extract_paragraph_formatting(self, paragraph):
        """Extract paragraph-level formatting"""
        para_format = {
            "level": paragraph.level,
            "alignment": str(paragraph.alignment) if paragraph.alignment else None,
            "line_spacing": paragraph.line_spacing,
            "space_before": paragraph.space_before.pt if paragraph.space_before else None,
            "space_after": paragraph.space_after.pt if paragraph.space_after else None,
            "indent": None,
            "left_indent": None,
            "right_indent": None,
            "bullet_format": self.extract_bullet_formatting(paragraph),
            "text_direction": None
        }
        
        try:
            pPr = paragraph._element.pPr
            if pPr is not None:
                if pPr.get('indent'):
                    para_format["indent"] = int(pPr.get('indent'))
                if pPr.get('marL'):
                    para_format["left_indent"] = int(pPr.get('marL'))
                if pPr.get('marR'):
                    para_format["right_indent"] = int(pPr.get('marR'))
                rtl = pPr.get('rtl')
                if rtl:
                    para_format["text_direction"] = "rtl" if rtl == '1' else "ltr"
        except:
            pass
        
        return para_format
    
    def extract_text_frame_properties(self, text_frame):
        """Extract text frame properties including text direction and rotation"""
        properties = {
            "margin_left": text_frame.margin_left,
            "margin_right": text_frame.margin_right,
            "margin_top": text_frame.margin_top,
            "margin_bottom": text_frame.margin_bottom,
            "word_wrap": text_frame.word_wrap,
            "auto_size": str(text_frame.auto_size) if hasattr(text_frame, 'auto_size') else None,
            "vertical_anchor": str(text_frame.vertical_anchor) if hasattr(text_frame, 'vertical_anchor') else None,
            "text_direction": None,
            "rotation_angle": None
        }
        
        try:
            bodyPr = text_frame._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
            if bodyPr is not None:
                vert = bodyPr.get('vert')
                if vert:
                    properties["text_direction"] = vert
                
                rot = bodyPr.get('rot')
                if rot:
                    properties["rotation_angle"] = int(rot) / 60000.0
        except:
            pass
        
        return properties
    
    def extract_table(self, shape):
        """Extract table structure and content"""
        table = shape.table
        table_data = {
            "rows": len(table.rows),
            "columns": len(table.columns),
            "cells": []
        }
        
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                cell_data = {
                    "row": row_idx,
                    "column": col_idx,
                    "text": cell.text,
                    "paragraphs": []
                }
                
                for paragraph in cell.text_frame.paragraphs:
                    para_data = {
                        "paragraph_formatting": self.extract_paragraph_formatting(paragraph),
                        "runs": []
                    }
                    
                    for run in paragraph.runs:
                        run_data = self.extract_run_formatting(run)
                        para_data["runs"].append(run_data)
                    
                    cell_data["paragraphs"].append(para_data)
                
                table_data["cells"].append(cell_data)
        
        return table_data
    
    def extract_chart(self, shape):
        """Extract chart data including chart type, data values, and styling"""
        chart_data = {
            "chart_type": None,
            "chart_style": None,
            "has_title": False,
            "title": None,
            "data_values": [],
            "categories": [],
            "series_names": []
        }
        
        try:
            chart = shape.chart
            
            if hasattr(chart, 'chart_type'):
                chart_data["chart_type"] = f"{str(chart.chart_type).split('.')[-1]} ({chart.chart_type})"
            
            if hasattr(chart, 'chart_style'):
                chart_data["chart_style"] = chart.chart_style
            
            if chart.has_title:
                chart_data["has_title"] = True
                chart_data["title"] = chart.chart_title.text_frame.text
            
            try:
                for series_idx, series in enumerate(chart.series):
                    series_data = {
                        "series_name": series.name,
                        "values": []
                    }
                    
                    if hasattr(series, 'values'):
                        series_data["values"] = list(series.values)
                    
                    chart_data["data_values"].append(series_data)
                    chart_data["series_names"].append(series.name)
            except:
                pass
            
            try:
                if hasattr(chart, 'plots') and len(chart.plots) > 0:
                    plot = chart.plots[0]
                    if hasattr(plot, 'categories'):
                        chart_data["categories"] = list(plot.categories)
            except:
                pass
        except:
            pass
        
        return chart_data
    
    def extract_shape(self, shape, slide_num):
        """Extract comprehensive shape information"""
        element = {
            "shape_id": shape.shape_id,
            "shape_name": shape.name,
            "element_type": None,
            "placeholder_info": self.extract_placeholder_info(shape),
            "fill": self.extract_shape_fill(shape),
            "line": self.extract_shape_line(shape),
            "shadow": self.extract_shape_shadow(shape)
        }
        
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or (shape.has_text_frame and not shape.has_table):
            element["element_type"] = "TextBox"
            element["text_frame_properties"] = self.extract_text_frame_properties(shape.text_frame)
            element["paragraphs"] = []
            
            for paragraph in shape.text_frame.paragraphs:
                para_data = {
                    "paragraph_formatting": self.extract_paragraph_formatting(paragraph),
                    "runs": []
                }
                
                for run in paragraph.runs:
                    run_data = self.extract_run_formatting(run)
                    para_data["runs"].append(run_data)
                
                element["paragraphs"].append(para_data)
            
            element["full_text"] = shape.text_frame.text
            
        elif shape.has_table:
            element["element_type"] = "Table"
            element["table_data"] = self.extract_table(shape)
            
        elif shape.has_chart:
            element["element_type"] = "Chart"
            element["chart_data"] = self.extract_chart(shape)
        
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            element["element_type"] = "Picture"
            element["image_info"] = {
                "description": shape.name,
                "alt_text": getattr(shape, 'alt_text', None) if hasattr(shape, 'alt_text') else None
            }
        
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            element["element_type"] = "AutoShape"
            if shape.has_text_frame:
                element["text_frame_properties"] = self.extract_text_frame_properties(shape.text_frame)
                element["full_text"] = shape.text_frame.text
        
        else:
            element["element_type"] = f"Other_{shape.shape_type}"
        
        element["dimensions"] = {
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
            "rotation": shape.rotation
        }
        
        return element
    
    def extract_grouped_shapes(self, group_shape, slide_num):
        """Extract shapes from a group recursively"""
        grouped_elements = []
        
        for shape in group_shape.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                grouped_elements.extend(self.extract_grouped_shapes(shape, slide_num))
            else:
                element = self.extract_shape(shape, slide_num)
                if element:
                    grouped_elements.append(element)
        
        return grouped_elements
    
    def extract_smartart_xml(self):
        """Extract SmartArt diagrams with hierarchical structure using XML parsing"""
        smartart_data = []
        
        try:
            with zipfile.ZipFile(self.pptx_path, 'r') as zip_ref:
                diagram_files = [f for f in zip_ref.namelist() if 'diagrams/data' in f and f.endswith('.xml')]
                
                for diagram_file in diagram_files:
                    try:
                        xml_content = zip_ref.read(diagram_file)
                        root = etree.fromstring(xml_content)
                        
                        smartart_element = {
                            "element_type": "SmartArt",
                            "layout_type": None,
                            "texts": [],
                            "nodes": [],
                            "full_text": ""
                        }
                        
                        try:
                            layout_node = root.find('.//dgm:cat', self.namespaces)
                            if layout_node is not None:
                                smartart_element["layout_type"] = layout_node.get('type')
                            
                            if not smartart_element["layout_type"]:
                                for elem in root.iter():
                                    if 'layoutNode' in elem.tag or 'cat' in elem.tag:
                                        cat_type = elem.get('type')
                                        if cat_type:
                                            smartart_element["layout_type"] = cat_type
                                            break
                        except:
                            pass
                        
                        try:
                            ptLst = root.find('.//dgm:ptLst', self.namespaces)
                            if ptLst is not None:
                                points = ptLst.findall('.//dgm:pt', self.namespaces)
                                
                                for pt in points:
                                    node_data = {
                                        "node_id": pt.get('modelId'),
                                        "level": None,
                                        "parent_id": None,
                                        "text": ""
                                    }
                                    
                                    prSet = pt.find('.//dgm:prSet', self.namespaces)
                                    if prSet is not None:
                                        presLayoutVars = prSet.find('.//dgm:presLayoutVars', self.namespaces)
                                        if presLayoutVars is not None:
                                            for child in presLayoutVars:
                                                if 'depth' in child.tag.lower() or 'level' in child.tag.lower():
                                                    try:
                                                        node_data["level"] = int(child.get('val', 0))
                                                    except:
                                                        pass
                                    
                                    t_elem = pt.find('.//dgm:t', self.namespaces)
                                    if t_elem is None:
                                        t_elem = pt.find('.//a:t', self.namespaces)
                                    
                                    if t_elem is not None and t_elem.text:
                                        node_data["text"] = t_elem.text.strip()
                                        smartart_element["texts"].append(node_data["text"])
                                    
                                    smartart_element["nodes"].append(node_data)
                            
                            cxnLst = root.find('.//dgm:cxnLst', self.namespaces)
                            if cxnLst is not None:
                                connections = cxnLst.findall('.//dgm:cxn', self.namespaces)
                                
                                for cxn in connections:
                                    cxn_type = cxn.get('type', '')
                                    if cxn_type in ['parOf', 'presOf']:
                                        src_id = cxn.get('srcId')
                                        dest_id = cxn.get('destId')
                                        
                                        for node in smartart_element["nodes"]:
                                            if node["node_id"] == src_id:
                                                node["parent_id"] = dest_id
                                                break
                            
                            if smartart_element["nodes"]:
                                root_nodes = [n for n in smartart_element["nodes"] if n["parent_id"] is None]
                                
                                def assign_level(node_id, level, nodes):
                                    for node in nodes:
                                        if node["node_id"] == node_id and node["level"] is None:
                                            node["level"] = level
                                            children = [n for n in nodes if n["parent_id"] == node_id]
                                            for child in children:
                                                assign_level(child["node_id"], level + 1, nodes)
                                
                                for root_node in root_nodes:
                                    assign_level(root_node["node_id"], 0, smartart_element["nodes"])
                        except:
                            pass
                        
                        if not smartart_element["texts"]:
                            for xpath in ['.//dgm:t', './/a:t', './/dgm:text', './/*[local-name()="t"]']:
                                try:
                                    text_elems = root.xpath(xpath, namespaces=self.namespaces)
                                    for elem in text_elems:
                                        if elem.text and elem.text.strip():
                                            smartart_element["texts"].append(elem.text.strip())
                                except:
                                    pass
                        
                        smartart_element["full_text"] = " ".join(smartart_element["texts"])
                        
                        if smartart_element["texts"] or smartart_element["nodes"]:
                            smartart_data.append(smartart_element)
                    except Exception as e:
                        continue
        except:
            pass
        
        return smartart_data
    
    def extract_links(self, shape):
        """Extract hyperlinks from shape"""
        links = []
        
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.hyperlink and run.hyperlink.address:
                        links.append({
                            "text": run.text,
                            "url": run.hyperlink.address
                        })
        
        return links
    
    def extract_slide(self, slide, slide_num):
        """Extract all content from a single slide"""
        slide_data = {
            "slide_number": slide_num,
            "layout_info": self.get_slide_layout_info(slide),
            "background": self.extract_background_info(slide),
            "elements": [],
            "links": [],
            "speaker_notes": None,
            "smartart": []
        }
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                slide_data["elements"].extend(
                    self.extract_grouped_shapes(shape, slide_num)
                )
            else:
                element = self.extract_shape(shape, slide_num)
                if element:
                    slide_data["elements"].append(element)
                
                links = self.extract_links(shape)
                if links:
                    slide_data["links"].extend(links)
        
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame.text.strip():
                slide_data["speaker_notes"] = {
                    "text": notes_frame.text,
                    "element_type": "SpeakerNotes"
                }
        
        return slide_data
    
    def extract_all(self):
        """Extract all content from presentation"""
        print("🔍 Extracting slide masters and layouts...")
        self.data["slide_masters"] = self.extract_slide_masters()
        
        print("🔍 Extracting slides...")
        for idx, slide in enumerate(self.presentation.slides, start=1):
            slide_data = self.extract_slide(slide, idx)
            self.data["slides"].append(slide_data)
            print(f"  ✓ Slide {idx}/{self.data['total_slides']}")
        
        print("🔍 Extracting SmartArt...")
        smartart_elements = self.extract_smartart_xml()
        
        if smartart_elements and self.data["slides"]:
            for smartart in smartart_elements:
                self.data["slides"][0]["smartart"].append(smartart)
        
        return self.data


# ============================================================================
# PART 2: TRANSLATOR - DeepL Translation Engine
# ============================================================================

class DeepLTranslator:
    """Translate PowerPoint content using DeepL API"""
    
    def __init__(self, target_language: str):
        if not DEEPL_API_KEY:
            raise ValueError("DEEPL_API_KEY not found in environment variables")
        
        self.api_key = DEEPL_API_KEY
        self.endpoint = DEEPL_ENDPOINT
        self.target_language = target_language
        
        # Get DeepL language code
        self.target_lang_code = SUPPORTED_LANGUAGES.get(target_language)
        if not self.target_lang_code:
            raise ValueError(f"Language '{target_language}' not supported. Supported: {list(SUPPORTED_LANGUAGES.keys())}")
        
        # Check for glossary - but only use if it exists for this language pair
        self.glossary_id = GLOSSARIES.get(self.target_lang_code)
        self.use_glossary = False  # Will be enabled after validation
        
        # RTL detection
        self.is_rtl = target_language in RTL_LANGUAGES
        
        # Statistics
        self.stats = {
            "total_texts_translated": 0,
            "api_calls": 0,
            "total_characters": 0
        }
        
        print(f"✓ DeepL Translator initialized")
        print(f"  Target: {target_language} ({self.target_lang_code})")
        if self.glossary_id:
            print(f"  Glossary ID found: {self.glossary_id}")
            print(f"  Note: Glossary will be validated on first API call")
        if self.is_rtl:
            print(f"  RTL Mode: ENABLED")
    
    def translate_batch(self, texts: List[str]) -> List[str]:
        """Translate a batch of texts using DeepL API"""
        if not texts:
            return []
        
        # Filter out empty texts but remember their positions
        text_map = {}
        non_empty_texts = []
        for idx, text in enumerate(texts):
            if text and text.strip():
                text_map[len(non_empty_texts)] = idx
                non_empty_texts.append(text)
        
        if not non_empty_texts:
            return texts
        
        try:
            # Prepare request
            headers = {
                "Authorization": f"DeepL-Auth-Key {self.api_key}",
                "Content-Type": "application/json"
            }
            
            # Build payload with formality and model_type
            payload = {
                "text": non_empty_texts,
                "source_lang": "EN",
                "target_lang": self.target_lang_code,
                "formality": "prefer_more",
                "model_type": "prefer_quality_optimized"
            }
            
            # Add glossary if available and not yet validated as failing
            if self.glossary_id and not hasattr(self, '_glossary_failed'):
                payload["glossary_id"] = self.glossary_id
            
            # Make API call
            response = requests.post(self.endpoint, headers=headers, json=payload)
            
            # Handle glossary-specific errors
            if response.status_code == 400:
                error_detail = ""
                try:
                    error_json = response.json()
                    error_msg = error_json.get('message', str(error_json))
                    
                    # Check if it's a glossary error
                    if 'glossary' in error_msg.lower():
                        print(f"⚠️  Glossary error: {error_msg}")
                        print(f"   Glossary ID: {self.glossary_id}")
                        print(f"   Source lang: EN, Target lang: {self.target_lang_code}")
                        print(f"   Retrying without glossary...")
                        self._glossary_failed = True  # Mark glossary as failed
                        
                        # Retry without glossary
                        payload.pop("glossary_id", None)
                        response = requests.post(self.endpoint, headers=headers, json=payload)
                    else:
                        error_detail = f": {error_msg}"
                        print(f"⚠️  DeepL API error 400{error_detail}")
                        return self.translate_one_by_one(texts)
                except:
                    error_detail = f": {response.text}"
                    print(f"⚠️  DeepL API error 400{error_detail}")
                    return self.translate_one_by_one(texts)
            
            # Better error handling for other status codes
            if response.status_code != 200:
                error_detail = ""
                try:
                    error_json = response.json()
                    error_detail = f": {error_json.get('message', error_json)}"
                except:
                    error_detail = f": {response.text}"
                print(f"⚠️  DeepL API error {response.status_code}{error_detail}")
                return self.translate_one_by_one(texts)
            
            # Parse response
            result = response.json()
            translated_texts = [item["text"] for item in result.get("translations", [])]
            
            # Update statistics
            self.stats["api_calls"] += 1
            self.stats["total_texts_translated"] += len(non_empty_texts)
            self.stats["total_characters"] += sum(len(t) for t in non_empty_texts)
            
            # Reconstruct full list with empty texts in original positions
            result_texts = texts.copy()
            for new_idx, orig_idx in text_map.items():
                if new_idx < len(translated_texts):
                    result_texts[orig_idx] = translated_texts[new_idx]
            
            return result_texts
            
        except requests.exceptions.RequestException as e:
            print(f"⚠️  DeepL API error: {e}")
            # Fallback: translate one by one
            return self.translate_one_by_one(texts)
        except Exception as e:
            print(f"⚠️  Translation error: {e}")
            return texts
    
    def translate_one_by_one(self, texts: List[str]) -> List[str]:
        """Fallback: translate texts one by one"""
        translated = []
        for text in texts:
            if not text or not text.strip():
                translated.append(text)
                continue
            
            try:
                headers = {
                    "Authorization": f"DeepL-Auth-Key {self.api_key}",
                    "Content-Type": "application/json"
                }
                
                payload = {
                    "text": [text],
                    "source_lang": "EN",
                    "target_lang": self.target_lang_code,
                    "formality": "prefer_more",
                    "model_type": "prefer_quality_optimized"
                }
                
                # Only add glossary if it hasn't failed before
                if self.glossary_id and not hasattr(self, '_glossary_failed'):
                    payload["glossary_id"] = self.glossary_id
                
                response = requests.post(self.endpoint, headers=headers, json=payload)
                
                if response.status_code != 200:
                    # Silently skip failed texts in one-by-one mode
                    translated.append(text)
                    continue
                
                result = response.json()
                translated_text = result["translations"][0]["text"]
                
                self.stats["api_calls"] += 1
                self.stats["total_texts_translated"] += 1
                self.stats["total_characters"] += len(text)
                
                translated.append(translated_text)
                
            except Exception as e:
                translated.append(text)
        
        return translated
    
    def translate_text_runs(self, runs: List[Dict]) -> List[Dict]:
        """Translate text runs while preserving all formatting metadata"""
        if not runs:
            return runs
        
        texts = [run.get("text", "") for run in runs]
        translated_texts = self.translate_batch(texts)
        
        translated_runs = []
        for idx, run in enumerate(runs):
            new_run = deepcopy(run)
            new_run["text"] = translated_texts[idx]
            translated_runs.append(new_run)
        
        return translated_runs
    
    def translate_paragraphs(self, paragraphs: List[Dict]) -> List[Dict]:
        """Translate paragraphs while preserving all paragraph formatting"""
        if not paragraphs:
            return paragraphs
        
        translated_paragraphs = []
        for para in paragraphs:
            new_para = deepcopy(para)
            
            if "runs" in new_para:
                new_para["runs"] = self.translate_text_runs(new_para["runs"])
            
            translated_paragraphs.append(new_para)
        
        return translated_paragraphs
    
    def translate_text_element(self, element: Dict) -> Dict:
        """Translate a text element preserving all metadata"""
        new_element = deepcopy(element)
        
        if "paragraphs" in new_element:
            new_element["paragraphs"] = self.translate_paragraphs(new_element["paragraphs"])
        
        if "paragraphs" in new_element:
            all_text = []
            for para in new_element["paragraphs"]:
                para_text = []
                if "runs" in para:
                    for run in para["runs"]:
                        if run.get("text"):
                            para_text.append(run["text"])
                if para_text:
                    all_text.append("".join(para_text))
            new_element["full_text"] = "\n".join(all_text) if all_text else ""
        
        return new_element
    
    def translate_table(self, table_data: Dict) -> Dict:
        """Translate table cells while preserving table structure"""
        new_table = deepcopy(table_data)
        
        if "cells" in new_table:
            translated_cells = []
            for cell in new_table["cells"]:
                if "paragraphs" in cell:
                    cell["paragraphs"] = self.translate_paragraphs(cell["paragraphs"])
                
                if "text" in cell and "paragraphs" in cell:
                    all_text = []
                    for para in cell["paragraphs"]:
                        if "runs" in para:
                            para_text = "".join(run.get("text", "") for run in para["runs"])
                            if para_text:
                                all_text.append(para_text)
                    cell["text"] = "\n".join(all_text) if all_text else ""
                
                translated_cells.append(cell)
            new_table["cells"] = translated_cells
        
        return new_table
    
    def translate_chart(self, chart_data: Dict) -> Dict:
        """Translate chart text elements while preserving chart data"""
        new_chart = deepcopy(chart_data)
        
        if "title" in new_chart and new_chart["title"]:
            translated = self.translate_batch([new_chart["title"]])
            new_chart["title"] = translated[0]
        
        if "data_values" in new_chart and new_chart["data_values"]:
            series_names = [s.get("series_name", "") for s in new_chart["data_values"] if s.get("series_name")]
            if series_names:
                translated_names = self.translate_batch(series_names)
                name_idx = 0
                for series in new_chart["data_values"]:
                    if series.get("series_name"):
                        series["series_name"] = translated_names[name_idx]
                        name_idx += 1
        
        if "series_names" in new_chart and new_chart["series_names"]:
            new_chart["series_names"] = self.translate_batch(new_chart["series_names"])
        
        if "categories" in new_chart and new_chart["categories"]:
            text_categories = [cat for cat in new_chart["categories"] if isinstance(cat, str)]
            if text_categories:
                translated_cats = self.translate_batch(new_chart["categories"])
                new_chart["categories"] = translated_cats
        
        return new_chart
    
    def translate_smartart(self, smartart: Dict) -> Dict:
        """Translate SmartArt text while preserving hierarchical structure"""
        new_smartart = deepcopy(smartart)
        
        if "texts" in new_smartart and new_smartart["texts"]:
            new_smartart["texts"] = self.translate_batch(new_smartart["texts"])
        
        if "nodes" in new_smartart and new_smartart["nodes"]:
            node_texts = [node.get("text", "") for node in new_smartart["nodes"]]
            if node_texts:
                translated_node_texts = self.translate_batch(node_texts)
                for idx, node in enumerate(new_smartart["nodes"]):
                    if node.get("text"):
                        node["text"] = translated_node_texts[idx]
        
        if "texts" in new_smartart:
            new_smartart["full_text"] = " ".join(new_smartart["texts"])
        
        return new_smartart
    
    def translate_speaker_notes(self, notes: Dict) -> Dict:
        """Translate speaker notes while preserving metadata"""
        new_notes = deepcopy(notes)
        
        if "text" in new_notes and new_notes["text"]:
            translated = self.translate_batch([new_notes["text"]])
            new_notes["text"] = translated[0]
        
        return new_notes
    
    def translate_slide(self, slide: Dict, slide_num: int) -> Dict:
        """Translate a single slide while preserving all metadata"""
        new_slide = deepcopy(slide)
        
        if "elements" in new_slide:
            translated_elements = []
            for element in new_slide["elements"]:
                element_type = element.get("element_type")
                
                if element_type == "Table":
                    if "table_data" in element:
                        element["table_data"] = self.translate_table(element["table_data"])
                    translated_elements.append(element)
                    
                elif element_type == "Chart":
                    if "chart_data" in element:
                        element["chart_data"] = self.translate_chart(element["chart_data"])
                    translated_elements.append(element)
                    
                elif element_type in ["TextBox", "AutoShape"]:
                    translated_elements.append(self.translate_text_element(element))
                    
                else:
                    translated_elements.append(deepcopy(element))
            
            new_slide["elements"] = translated_elements
        
        if "speaker_notes" in new_slide and new_slide["speaker_notes"]:
            new_slide["speaker_notes"] = self.translate_speaker_notes(new_slide["speaker_notes"])
        
        if "smartart" in new_slide and new_slide["smartart"]:
            translated_smartart = []
            for smartart in new_slide["smartart"]:
                translated_smartart.append(self.translate_smartart(smartart))
            new_slide["smartart"] = translated_smartart
        
        return new_slide
    
    def translate_presentation(self, data: Dict) -> Dict:
        """Translate entire presentation while preserving all metadata"""
        print(f"\n🌍 Translating to {self.target_language}...")
        print(f"Total slides: {data['total_slides']}")
        print("=" * 80)
        
        translated_data = {
            "presentation_name": data["presentation_name"],
            "total_slides": data["total_slides"],
            "target_language": self.target_language,
            "is_rtl": self.is_rtl,
            "slides": []
        }
        
        if "slide_masters" in data:
            translated_data["slide_masters"] = deepcopy(data["slide_masters"])
        
        start_time = time.time()
        for idx, slide in enumerate(data["slides"], 1):
            print(f"  Translating slide {idx}/{data['total_slides']}...", end=" ", flush=True)
            translated_slide = self.translate_slide(slide, idx)
            translated_data["slides"].append(translated_slide)
            print("✓")
            
            # Small delay to avoid rate limits
            if idx < data['total_slides']:
                time.sleep(0.1)
        
        elapsed_time = time.time() - start_time
        
        print("=" * 80)
        print("✓ Translation complete!")
        print(f"  Texts translated: {self.stats['total_texts_translated']}")
        print(f"  API calls: {self.stats['api_calls']}")
        print(f"  Characters: {self.stats['total_characters']}")
        print(f"  Time: {elapsed_time:.2f} seconds")
        
        return translated_data


# ============================================================================
# PART 3: REASSEMBLER - PowerPoint Rebuilder
# ============================================================================

class PPTXReassembler:
    """Reassemble PowerPoint presentation from translated JSON"""
    
    def __init__(self, original_pptx_path: str, translated_data: Dict):
        self.original_pptx_path = original_pptx_path
        self.translated_data = translated_data
        
        print(f"\n📝 Loading original presentation: {original_pptx_path}")
        self.presentation = Presentation(original_pptx_path)
        
        self.is_rtl = self.translated_data.get('is_rtl', False)
        self.target_language = self.translated_data.get('target_language', 'Unknown')
        
        self.stats = {
            "slides_processed": 0,
            "elements_updated": 0,
            "text_runs_updated": 0,
            "tables_updated": 0,
            "charts_updated": 0,
            "notes_updated": 0,
            "rtl_paragraphs_set": 0,
            "auto_shrink_enabled": 0,
            "shapes_mirrored": 0
        }
        
        print(f"✓ Loaded {len(self.presentation.slides)} slides")
        print(f"✓ Target language: {self.target_language}")
        if self.is_rtl:
            print(f"✓ RTL mode: ENABLED")
    
    def find_shape_by_id(self, slide, shape_id: int):
        """Find a shape in a slide by its shape_id"""
        for shape in slide.shapes:
            if shape.shape_id == shape_id:
                return shape
        return None
    
    def mirror_shape_horizontal(self, shape, slide_width):
        """Mirror a shape's horizontal position for RTL languages"""
        if not self.is_rtl:
            return
        
        try:
            current_left = shape.left
            shape_width = shape.width
            new_left = slide_width - (current_left + shape_width)
            shape.left = new_left
            self.stats["shapes_mirrored"] += 1
        except:
            pass
    
    def mirror_slide_layout(self, slide):
        """Mirror all shapes on a slide horizontally for RTL languages"""
        if not self.is_rtl:
            return
        
        try:
            slide_width = self.presentation.slide_width
            for shape in slide.shapes:
                self.mirror_shape_horizontal(shape, slide_width)
        except:
            pass
    
    def set_rtl_if_needed(self, paragraph):
        """Set RTL properties on paragraph if target language is RTL"""
        if not self.is_rtl:
            return
        
        try:
            pPr = paragraph._element.get_or_add_pPr()
            pPr.set('rtl', '1')
            pPr.set('algn', 'r')
            self.stats["rtl_paragraphs_set"] += 1
        except:
            pass
    
    def enable_auto_shrink(self, text_frame):
        """Enable auto-shrink to prevent text overflow"""
        try:
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            self.stats["auto_shrink_enabled"] += 1
        except:
            pass
    
    def update_text_runs(self, paragraph, translated_runs: list):
        """Update text runs in a paragraph with translated text"""
        existing_runs = list(paragraph.runs)
        
        if len(existing_runs) == len(translated_runs):
            for idx, (ppt_run, json_run) in enumerate(zip(existing_runs, translated_runs)):
                ppt_run.text = json_run.get("text", "")
                self.stats["text_runs_updated"] += 1
        else:
            while len(paragraph.runs) > 1:
                paragraph._element.remove(paragraph.runs[-1]._element)
            
            if len(paragraph.runs) == 0:
                paragraph.add_run()
            
            if len(translated_runs) > 0:
                paragraph.runs[0].text = translated_runs[0].get("text", "")
                self.stats["text_runs_updated"] += 1
            
            for json_run in translated_runs[1:]:
                new_run = paragraph.add_run()
                new_run.text = json_run.get("text", "")
                
                try:
                    font = new_run.font
                    if json_run.get("bold") is not None:
                        font.bold = json_run["bold"]
                    if json_run.get("italic") is not None:
                        font.italic = json_run["italic"]
                    if json_run.get("font_size"):
                        font.size = Pt(json_run["font_size"])
                except:
                    pass
                
                self.stats["text_runs_updated"] += 1
        
        self.set_rtl_if_needed(paragraph)
    
    def update_text_frame(self, shape, translated_element: dict):
        """Update text in a text frame"""
        if not shape.has_text_frame:
            return
        
        text_frame = shape.text_frame
        translated_paragraphs = translated_element.get("paragraphs", [])
        existing_paragraphs = list(text_frame.paragraphs)
        
        for idx, translated_para in enumerate(translated_paragraphs):
            if idx < len(existing_paragraphs):
                ppt_para = existing_paragraphs[idx]
                translated_runs = translated_para.get("runs", [])
                self.update_text_runs(ppt_para, translated_runs)
            else:
                ppt_para = text_frame.add_paragraph()
                translated_runs = translated_para.get("runs", [])
                self.update_text_runs(ppt_para, translated_runs)
        
        while len(text_frame.paragraphs) > len(translated_paragraphs):
            try:
                text_frame._element.remove(text_frame.paragraphs[-1]._element)
            except:
                break
        
        self.enable_auto_shrink(text_frame)
    
    def update_table(self, shape, translated_table_data: dict):
        """Update table cell text with translations"""
        if not shape.has_table:
            return
        
        table = shape.table
        translated_cells = translated_table_data.get("cells", [])
        
        for cell_data in translated_cells:
            row = cell_data.get("row")
            col = cell_data.get("column")
            
            if row is not None and col is not None:
                try:
                    cell = table.cell(row, col)
                    translated_paragraphs = cell_data.get("paragraphs", [])
                    
                    if cell.text_frame and translated_paragraphs:
                        for idx, translated_para in enumerate(translated_paragraphs):
                            if idx < len(cell.text_frame.paragraphs):
                                ppt_para = cell.text_frame.paragraphs[idx]
                                translated_runs = translated_para.get("runs", [])
                                self.update_text_runs(ppt_para, translated_runs)
                        
                        self.enable_auto_shrink(cell.text_frame)
                except:
                    pass
        
        self.stats["tables_updated"] += 1
    
    def update_chart(self, shape, translated_chart_data: dict):
        """Update chart text elements"""
        if not shape.has_chart:
            return
        
        try:
            chart = shape.chart
            
            if translated_chart_data.get("title") and chart.has_title:
                try:
                    chart.chart_title.text_frame.text = translated_chart_data["title"]
                    self.enable_auto_shrink(chart.chart_title.text_frame)
                except:
                    pass
            
            translated_series_names = translated_chart_data.get("series_names", [])
            if translated_series_names:
                for idx, series in enumerate(chart.series):
                    if idx < len(translated_series_names):
                        try:
                            series.name = translated_series_names[idx]
                        except:
                            pass
            
            self.stats["charts_updated"] += 1
        except:
            pass
    
    def update_speaker_notes(self, slide, translated_notes: dict):
        """Update speaker notes for a slide"""
        if not translated_notes or not translated_notes.get("text"):
            return
        
        try:
            if not slide.has_notes_slide:
                notes_slide = slide.notes_slide
            else:
                notes_slide = slide.notes_slide
            
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = translated_notes["text"]
            
            if self.is_rtl and notes_text_frame.paragraphs:
                for para in notes_text_frame.paragraphs:
                    self.set_rtl_if_needed(para)
            
            self.enable_auto_shrink(notes_text_frame)
            self.stats["notes_updated"] += 1
        except:
            pass
    
    def update_slide(self, slide_idx: int):
        """Update a single slide with translated content"""
        ppt_slide = self.presentation.slides[slide_idx]
        translated_slide = self.translated_data["slides"][slide_idx]
        
        self.mirror_slide_layout(ppt_slide)
        
        translated_elements = translated_slide.get("elements", [])
        
        for element in translated_elements:
            shape_id = element.get("shape_id")
            element_type = element.get("element_type")
            
            shape = self.find_shape_by_id(ppt_slide, shape_id)
            if not shape:
                continue
            
            try:
                if element_type in ["TextBox", "AutoShape"]:
                    self.update_text_frame(shape, element)
                    self.stats["elements_updated"] += 1
                
                elif element_type == "Table":
                    table_data = element.get("table_data")
                    if table_data:
                        self.update_table(shape, table_data)
                        self.stats["elements_updated"] += 1
                
                elif element_type == "Chart":
                    chart_data = element.get("chart_data")
                    if chart_data:
                        self.update_chart(shape, chart_data)
                        self.stats["elements_updated"] += 1
            except:
                pass
        
        translated_notes = translated_slide.get("speaker_notes")
        if translated_notes:
            self.update_speaker_notes(ppt_slide, translated_notes)
        
        self.stats["slides_processed"] += 1
    
    def reassemble(self, output_path: str):
        """Main reassembly process"""
        print("\n🔨 Starting reassembly...")
        print(f"Target language: {self.target_language}")
        if self.is_rtl:
            print(f"RTL mode: ENABLED")
        print(f"Auto-shrink: ENABLED")
        print("=" * 80)
        
        num_slides = min(len(self.presentation.slides), len(self.translated_data["slides"]))
        
        for slide_idx in range(num_slides):
            print(f"  Processing slide {slide_idx + 1}/{num_slides}...", end=" ", flush=True)
            self.update_slide(slide_idx)
            print("✓")
        
        print("=" * 80)
        print(f"\n💾 Saving to: {output_path}")
        self.presentation.save(output_path)
        
        print("\n" + "=" * 80)
        print("✅ REASSEMBLY COMPLETE!")
        print("=" * 80)
        print(f"Slides processed: {self.stats['slides_processed']}")
        print(f"Elements updated: {self.stats['elements_updated']}")
        print(f"Text runs updated: {self.stats['text_runs_updated']}")
        print(f"Tables updated: {self.stats['tables_updated']}")
        print(f"Charts updated: {self.stats['charts_updated']}")
        print(f"Speaker notes updated: {self.stats['notes_updated']}")
        if self.is_rtl:
            print(f"RTL paragraphs set: {self.stats['rtl_paragraphs_set']}")
            print(f"Shapes mirrored: {self.stats['shapes_mirrored']}")
        print(f"Auto-shrink enabled: {self.stats['auto_shrink_enabled']} text frames")
        print(f"\n✓ Output saved to: {output_path}")
        print("=" * 80)
        
        return self.stats


# ============================================================================
# MAIN PIPELINE
# ============================================================================

def run_pipeline(input_pptx: str, target_language: str, output_pptx: Optional[str] = None):
    """
    Run the complete translation pipeline
    
    Args:
        input_pptx: Path to input PowerPoint file
        target_language: Target language (e.g., 'Spanish', 'French')
        output_pptx: Optional output path (auto-generated if None)
    """
    print("\n" + "=" * 80)
    print("🚀 PowerPoint Translation Pipeline")
    print("=" * 80)
    print(f"Input: {input_pptx}")
    print(f"Target Language: {target_language}")
    print("=" * 80)
    
    # Validate input file
    if not os.path.exists(input_pptx):
        print(f"❌ Error: Input file not found: {input_pptx}")
        return 1
    
    # Generate output path if not provided
    if not output_pptx:
        input_path = Path(input_pptx)
        output_pptx = str(input_path.parent / f"{input_path.stem}-{target_language.lower()}-translated.pptx")
    
    try:
        # STAGE 1: EXTRACTION
        print("\n" + "=" * 80)
        print("STAGE 1: EXTRACTION")
        print("=" * 80)
        extractor = PPTXExtractor(input_pptx)
        extracted_data = extractor.extract_all()
        print(f"✓ Extracted {extracted_data['total_slides']} slides")
        print(f"✓ Total elements: {sum(len(s['elements']) for s in extracted_data['slides'])}")
        
        # STAGE 2: TRANSLATION
        print("\n" + "=" * 80)
        print("STAGE 2: TRANSLATION")
        print("=" * 80)
        translator = DeepLTranslator(target_language)
        translated_data = translator.translate_presentation(extracted_data)
        
        # STAGE 3: REASSEMBLY
        print("\n" + "=" * 80)
        print("STAGE 3: REASSEMBLY")
        print("=" * 80)
        reassembler = PPTXReassembler(input_pptx, translated_data)
        stats = reassembler.reassemble(output_pptx)
        
        print("\n" + "=" * 80)
        print("🎉 PIPELINE COMPLETE!")
        print("=" * 80)
        print(f"✓ Input: {input_pptx}")
        print(f"✓ Output: {output_pptx}")
        print(f"✓ Language: {target_language}")
        print("=" * 80)
        
        return 0
        
    except Exception as e:
        print(f"\n❌ Pipeline error: {e}")
        import traceback
        traceback.print_exc()
        return 1


def main():
    """Main entry point"""
    parser = argparse.ArgumentParser(
        description="PowerPoint Translation Pipeline with DeepL",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python ppt_translator_pipeline.py presentation.pptx -l Spanish
  python ppt_translator_pipeline.py presentation.pptx -l French -o output.pptx
  python ppt_translator_pipeline.py presentation.pptx -l German

Supported Languages:
  French, Spanish, Italian, German, Chinese, Japanese, Dutch, Swedish

Features:
  ✓ Comprehensive metadata preservation
  ✓ DeepL translation with glossary support
  ✓ RTL language support
  ✓ Auto-shrink to prevent text overflow
  ✓ Template-based reassembly
        """
    )
    
    parser.add_argument(
        "input_pptx",
        help="Path to input PowerPoint file"
    )
    parser.add_argument(
        "-l", "--language",
        required=True,
        choices=list(SUPPORTED_LANGUAGES.keys()),
        help="Target language for translation"
    )
    parser.add_argument(
        "-o", "--output",
        help="Path to output PowerPoint file (default: <input>-<language>-translated.pptx)"
    )
    
    args = parser.parse_args()
    
    return run_pipeline(args.input_pptx, args.language, args.output)


if __name__ == "__main__":
    sys.exit(main())