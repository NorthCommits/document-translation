# import json
# import os
# from pptx import Presentation
# from pptx.util import Pt, Inches
# from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
# from copy import deepcopy
# import argparse

# class PPTXReassembler:
#     """
#     Reassembles PowerPoint presentation from translated JSON.
    
#     Strategy: Template-based approach
#     - Uses original PowerPoint as template
#     - Matches slides by layout information
#     - Replaces only text content
#     - Preserves all visual formatting automatically
    
#     This ensures 100% visual fidelity while applying translations.
#     """
    
#     def __init__(self, original_pptx_path: str, translated_json_path: str):
#         """
#         Initialize the reassembler.
        
#         Args:
#             original_pptx_path: Path to original PowerPoint file (template)
#             translated_json_path: Path to translated JSON file
#         """
#         self.original_pptx_path = original_pptx_path
#         self.translated_json_path = translated_json_path
        
#         # Load original presentation
#         print(f"Loading original presentation: {original_pptx_path}")
#         self.presentation = Presentation(original_pptx_path)
        
#         # Load translated JSON
#         print(f"Loading translated content: {translated_json_path}")
#         with open(translated_json_path, 'r', encoding='utf-8') as f:
#             self.translated_data = json.load(f)
        
#         # Statistics
#         self.stats = {
#             "slides_processed": 0,
#             "elements_updated": 0,
#             "text_runs_updated": 0,
#             "tables_updated": 0,
#             "charts_updated": 0,
#             "notes_updated": 0
#         }
        
#         print(f"✓ Loaded {len(self.presentation.slides)} slides from original")
#         print(f"✓ Loaded {len(self.translated_data['slides'])} slides from JSON")
    
#     def find_shape_by_id(self, slide, shape_id: int):
#         """
#         Find a shape in a slide by its shape_id.
        
#         Args:
#             slide: PowerPoint slide object
#             shape_id: Shape ID to find
            
#         Returns:
#             Shape object or None if not found
#         """
#         for shape in slide.shapes:
#             if shape.shape_id == shape_id:
#                 return shape
#         return None
    
#     def update_text_runs(self, paragraph, translated_runs: list):
#         """
#         Update text runs in a paragraph with translated text.
#         Preserves all formatting.
        
#         Args:
#             paragraph: PowerPoint paragraph object
#             translated_runs: List of translated run dictionaries
#         """
#         # Get existing runs
#         existing_runs = list(paragraph.runs)
        
#         # If counts match, update in place (safest)
#         if len(existing_runs) == len(translated_runs):
#             for idx, (ppt_run, json_run) in enumerate(zip(existing_runs, translated_runs)):
#                 ppt_run.text = json_run.get("text", "")
#                 self.stats["text_runs_updated"] += 1
#         else:
#             # Counts don't match - need to rebuild runs
#             # Clear all runs except first
#             while len(paragraph.runs) > 1:
#                 paragraph._element.remove(paragraph.runs[-1]._element)
            
#             # Update first run or create if none
#             if len(paragraph.runs) == 0:
#                 paragraph.add_run()
            
#             # Set first run
#             if len(translated_runs) > 0:
#                 paragraph.runs[0].text = translated_runs[0].get("text", "")
#                 self.stats["text_runs_updated"] += 1
            
#             # Add remaining runs
#             for json_run in translated_runs[1:]:
#                 new_run = paragraph.add_run()
#                 new_run.text = json_run.get("text", "")
                
#                 # Try to apply formatting from JSON
#                 try:
#                     font = new_run.font
#                     if json_run.get("bold") is not None:
#                         font.bold = json_run["bold"]
#                     if json_run.get("italic") is not None:
#                         font.italic = json_run["italic"]
#                     if json_run.get("font_size"):
#                         font.size = Pt(json_run["font_size"])
#                 except:
#                     pass
                
#                 self.stats["text_runs_updated"] += 1
    
#     def update_text_frame(self, shape, translated_element: dict):
#         """
#         Update text in a text frame (TextBox, AutoShape, etc.).
        
#         Args:
#             shape: PowerPoint shape with text_frame
#             translated_element: Translated element dictionary from JSON
#         """
#         if not shape.has_text_frame:
#             return
        
#         text_frame = shape.text_frame
#         translated_paragraphs = translated_element.get("paragraphs", [])
        
#         # Get existing paragraphs
#         existing_paragraphs = list(text_frame.paragraphs)
        
#         # Update existing paragraphs
#         for idx, translated_para in enumerate(translated_paragraphs):
#             if idx < len(existing_paragraphs):
#                 # Update existing paragraph
#                 ppt_para = existing_paragraphs[idx]
#                 translated_runs = translated_para.get("runs", [])
#                 self.update_text_runs(ppt_para, translated_runs)
#             else:
#                 # Add new paragraph
#                 ppt_para = text_frame.add_paragraph()
#                 translated_runs = translated_para.get("runs", [])
#                 self.update_text_runs(ppt_para, translated_runs)
        
#         # Remove extra paragraphs if JSON has fewer
#         while len(text_frame.paragraphs) > len(translated_paragraphs):
#             # Remove last paragraph
#             try:
#                 text_frame._element.remove(text_frame.paragraphs[-1]._element)
#             except:
#                 break
    
#     def update_table(self, shape, translated_table_data: dict):
#         """
#         Update table cell text with translations.
        
#         Args:
#             shape: PowerPoint shape with table
#             translated_table_data: Translated table_data dictionary from JSON
#         """
#         if not shape.has_table:
#             return
        
#         table = shape.table
#         translated_cells = translated_table_data.get("cells", [])
        
#         # Update each cell
#         for cell_data in translated_cells:
#             row = cell_data.get("row")
#             col = cell_data.get("column")
            
#             if row is not None and col is not None:
#                 try:
#                     cell = table.cell(row, col)
#                     translated_paragraphs = cell_data.get("paragraphs", [])
                    
#                     # Update cell text frame
#                     if cell.text_frame and translated_paragraphs:
#                         for idx, translated_para in enumerate(translated_paragraphs):
#                             if idx < len(cell.text_frame.paragraphs):
#                                 ppt_para = cell.text_frame.paragraphs[idx]
#                                 translated_runs = translated_para.get("runs", [])
#                                 self.update_text_runs(ppt_para, translated_runs)
                
#                 except Exception as e:
#                     print(f"  Warning: Could not update cell ({row}, {col}): {e}")
        
#         self.stats["tables_updated"] += 1
    
#     def update_chart(self, shape, translated_chart_data: dict):
#         """
#         Update chart text elements (title, series names, categories).
#         Note: Chart data values are NOT translated, only text labels.
        
#         Args:
#             shape: PowerPoint shape with chart
#             translated_chart_data: Translated chart_data dictionary from JSON
#         """
#         if not shape.has_chart:
#             return
        
#         try:
#             chart = shape.chart
            
#             # Update chart title
#             if translated_chart_data.get("title") and chart.has_title:
#                 try:
#                     chart.chart_title.text_frame.text = translated_chart_data["title"]
#                 except:
#                     pass
            
#             # Update series names
#             translated_series_names = translated_chart_data.get("series_names", [])
#             if translated_series_names:
#                 for idx, series in enumerate(chart.series):
#                     if idx < len(translated_series_names):
#                         try:
#                             series.name = translated_series_names[idx]
#                         except:
#                             pass
            
#             # Update categories (if they are text)
#             translated_categories = translated_chart_data.get("categories", [])
#             if translated_categories:
#                 try:
#                     # This is tricky - categories might not be directly settable
#                     # We'll try but might fail on some chart types
#                     pass  # Categories are usually in the chart data, hard to update
#                 except:
#                     pass
            
#             self.stats["charts_updated"] += 1
            
#         except Exception as e:
#             print(f"  Warning: Could not update chart: {e}")
    
#     def update_speaker_notes(self, slide, translated_notes: dict):
#         """
#         Update speaker notes for a slide.
        
#         Args:
#             slide: PowerPoint slide object
#             translated_notes: Translated speaker_notes dictionary from JSON
#         """
#         if not translated_notes or not translated_notes.get("text"):
#             return
        
#         try:
#             # Ensure slide has notes
#             if not slide.has_notes_slide:
#                 notes_slide = slide.notes_slide  # This creates it
#             else:
#                 notes_slide = slide.notes_slide
            
#             # Update notes text
#             notes_text_frame = notes_slide.notes_text_frame
#             notes_text_frame.text = translated_notes["text"]
            
#             self.stats["notes_updated"] += 1
            
#         except Exception as e:
#             print(f"  Warning: Could not update speaker notes: {e}")
    
#     def update_slide(self, slide_idx: int):
#         """
#         Update a single slide with translated content.
        
#         Args:
#             slide_idx: 0-based slide index
#         """
#         # Get PowerPoint slide (0-based)
#         ppt_slide = self.presentation.slides[slide_idx]
        
#         # Get translated slide data (slides array is 0-based in JSON)
#         translated_slide = self.translated_data["slides"][slide_idx]
        
#         print(f"Processing slide {slide_idx + 1}/{len(self.presentation.slides)}...", end=" ")
        
#         # Update elements
#         translated_elements = translated_slide.get("elements", [])
        
#         for element in translated_elements:
#             shape_id = element.get("shape_id")
#             element_type = element.get("element_type")
            
#             # Find shape by ID
#             shape = self.find_shape_by_id(ppt_slide, shape_id)
            
#             if not shape:
#                 continue
            
#             # Update based on element type
#             try:
#                 if element_type in ["TextBox", "AutoShape"]:
#                     self.update_text_frame(shape, element)
#                     self.stats["elements_updated"] += 1
                
#                 elif element_type == "Table":
#                     table_data = element.get("table_data")
#                     if table_data:
#                         self.update_table(shape, table_data)
#                         self.stats["elements_updated"] += 1
                
#                 elif element_type == "Chart":
#                     chart_data = element.get("chart_data")
#                     if chart_data:
#                         self.update_chart(shape, chart_data)
#                         self.stats["elements_updated"] += 1
                
#                 # Picture and other types don't need text updates
                
#             except Exception as e:
#                 print(f"\n  Warning: Error updating shape {shape_id} ({element_type}): {e}")
        
#         # Update speaker notes
#         translated_notes = translated_slide.get("speaker_notes")
#         if translated_notes:
#             self.update_speaker_notes(ppt_slide, translated_notes)
        
#         self.stats["slides_processed"] += 1
#         print("✓")
    
#     def verify_slide_count(self):
#         """
#         Verify that original PPT and translated JSON have same number of slides.
        
#         Returns:
#             bool: True if counts match
#         """
#         ppt_count = len(self.presentation.slides)
#         json_count = len(self.translated_data["slides"])
        
#         if ppt_count != json_count:
#             print(f"\n⚠️  WARNING: Slide count mismatch!")
#             print(f"   Original PPT: {ppt_count} slides")
#             print(f"   Translated JSON: {json_count} slides")
#             print(f"   Will process minimum: {min(ppt_count, json_count)} slides")
#             return False
        
#         return True
    
#     def reassemble(self, output_path: str):
#         """
#         Main reassembly process.
        
#         Args:
#             output_path: Path to save the reassembled PowerPoint
#         """
#         print("\n" + "=" * 80)
#         print("STARTING REASSEMBLY")
#         print("=" * 80)
        
#         # Verify slide counts
#         self.verify_slide_count()
        
#         # Process each slide
#         num_slides = min(len(self.presentation.slides), len(self.translated_data["slides"]))
        
#         print(f"\nProcessing {num_slides} slides...")
#         print("-" * 80)
        
#         for slide_idx in range(num_slides):
#             self.update_slide(slide_idx)
        
#         # Save reassembled presentation
#         print("-" * 80)
#         print(f"\nSaving reassembled presentation to: {output_path}")
#         self.presentation.save(output_path)
        
#         # Print statistics
#         print("\n" + "=" * 80)
#         print("REASSEMBLY COMPLETE!")
#         print("=" * 80)
#         print(f"Slides processed: {self.stats['slides_processed']}")
#         print(f"Elements updated: {self.stats['elements_updated']}")
#         print(f"Text runs updated: {self.stats['text_runs_updated']}")
#         print(f"Tables updated: {self.stats['tables_updated']}")
#         print(f"Charts updated: {self.stats['charts_updated']}")
#         print(f"Speaker notes updated: {self.stats['notes_updated']}")
#         print(f"\n✓ Output saved to: {output_path}")
#         print("=" * 80)
        
#         return self.stats


# def main():
#     """Main function to run reassembly"""
#     parser = argparse.ArgumentParser(
#         description="Reassemble PowerPoint from translated JSON",
#         formatter_class=argparse.RawDescriptionHelpFormatter,
#         epilog="""
# Examples:
#   # Basic usage
#   python reassembler.py original.pptx translated.json output.pptx
  
#   # Using default output name
#   python reassembler.py original.pptx translated.json
  
#   # The output will be: translated_reassembled.pptx
#         """
#     )
    
#     parser.add_argument(
#         "original_pptx",
#         help="Path to original PowerPoint file (template)"
#     )
#     parser.add_argument(
#         "translated_json",
#         help="Path to translated JSON file"
#     )
#     parser.add_argument(
#         "output_pptx",
#         nargs='?',
#         help="Path to output PowerPoint file (optional, default: <translated_json>_reassembled.pptx)"
#     )
    
#     args = parser.parse_args()
    
#     # Determine output path
#     if args.output_pptx:
#         output_path = args.output_pptx
#     else:
#         # Generate default output name
#         json_base = args.translated_json.replace(".json", "")
#         output_path = f"{json_base}_reassembled.pptx"
    
#     # Check if files exist
#     if not os.path.exists(args.original_pptx):
#         print(f"Error: Original PowerPoint file not found: {args.original_pptx}")
#         return 1
    
#     if not os.path.exists(args.translated_json):
#         print(f"Error: Translated JSON file not found: {args.translated_json}")
#         return 1
    
#     # Create reassembler and run
#     try:
#         reassembler = PPTXReassembler(args.original_pptx, args.translated_json)
#         stats = reassembler.reassemble(output_path)
#         return 0
#     except Exception as e:
#         print(f"\n❌ Error during reassembly: {e}")
#         import traceback
#         traceback.print_exc()
#         return 1


# if __name__ == "__main__":
#     exit(main())



import json
import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from copy import deepcopy
import argparse

class PPTXReassembler:
    """
    Reassembles PowerPoint presentation from translated JSON.
    
    Strategy: Template-based approach
    - Uses original PowerPoint as template
    - Matches slides by layout information
    - Replaces only text content
    - Preserves all visual formatting automatically
    - Supports RTL (Right-to-Left) languages like Arabic, Hebrew
    - Prevents text overflow with auto-shrink
    
    This ensures 100% visual fidelity while applying translations.
    """
    
    def __init__(self, original_pptx_path: str, translated_json_path: str):
        """
        Initialize the reassembler.
        
        Args:
            original_pptx_path: Path to original PowerPoint file (template)
            translated_json_path: Path to translated JSON file
        """
        self.original_pptx_path = original_pptx_path
        self.translated_json_path = translated_json_path
        
        # Load original presentation
        print(f"Loading original presentation: {original_pptx_path}")
        self.presentation = Presentation(original_pptx_path)
        
        # Load translated JSON
        print(f"Loading translated content: {translated_json_path}")
        with open(translated_json_path, 'r', encoding='utf-8') as f:
            self.translated_data = json.load(f)
        
        # Check if target language is RTL
        self.is_rtl = self.translated_data.get('is_rtl', False)
        self.target_language = self.translated_data.get('target_language', 'Unknown')
        
        # Statistics
        self.stats = {
            "slides_processed": 0,
            "elements_updated": 0,
            "text_runs_updated": 0,
            "tables_updated": 0,
            "charts_updated": 0,
            "notes_updated": 0,
            "rtl_paragraphs_set": 0,
            "auto_shrink_enabled": 0,
            "shapes_mirrored": 0
        }
        
        print(f"✓ Loaded {len(self.presentation.slides)} slides from original")
        print(f"✓ Loaded {len(self.translated_data['slides'])} slides from JSON")
        print(f"✓ Target language: {self.target_language}")
        if self.is_rtl:
            print(f"✓ RTL mode: ENABLED (layout will be mirrored)")
    
    def find_shape_by_id(self, slide, shape_id: int):
        """
        Find a shape in a slide by its shape_id.
        Searches recursively through grouped shapes.
        
        Args:
            slide: PowerPoint slide object
            shape_id: Shape ID to find
            
        Returns:
            Shape object or None if not found
        """
        def search_shapes(shapes):
            """Recursively search through shapes including groups"""
            for shape in shapes:
                if shape.shape_id == shape_id:
                    return shape
                # If it's a group, search inside it
                if hasattr(shape, 'shape_type'):
                    try:
                        from pptx.enum.shapes import MSO_SHAPE_TYPE
                        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                            if hasattr(shape, 'shapes'):
                                found = search_shapes(shape.shapes)
                                if found:
                                    return found
                    except:
                        pass
            return None
        
        return search_shapes(slide.shapes)
    
    def mirror_shape_horizontal(self, shape, slide_width):
        """
        Mirror a shape's horizontal position (flip across the vertical center axis).
        This is used for RTL languages to flip the entire slide layout.
        
        Args:
            shape: PowerPoint shape object
            slide_width: Width of the slide in EMUs (English Metric Units)
        """
        if not self.is_rtl:
            return
        
        try:
            # Get current position
            current_left = shape.left
            shape_width = shape.width
            
            # Calculate mirrored position
            # Formula: new_left = slide_width - (current_left + shape_width)
            new_left = slide_width - (current_left + shape_width)
            
            # Set new position
            shape.left = new_left
            
            self.stats["shapes_mirrored"] += 1
            
        except Exception as e:
            # Some shapes might not support repositioning (like certain placeholders)
            pass
    
    def mirror_slide_layout(self, slide):
        """
        Mirror all shapes on a slide horizontally for RTL languages.
        This flips the entire slide layout so images move to the left
        and text areas move to the right.
        
        Args:
            slide: PowerPoint slide object
        """
        if not self.is_rtl:
            return
        
        try:
            # Get slide width
            slide_width = self.presentation.slide_width
            
            # Mirror each shape
            for shape in slide.shapes:
                self.mirror_shape_horizontal(shape, slide_width)
                
        except Exception as e:
            # If mirroring fails, continue without it
            pass
    
    def set_rtl_if_needed(self, paragraph):
        """
        Set RTL (right-to-left) properties on paragraph if target language is RTL.
        
        Args:
            paragraph: PowerPoint paragraph object
        """
        if not self.is_rtl:
            return
        
        try:
            # Access paragraph properties XML element
            pPr = paragraph._element.get_or_add_pPr()
            
            # Set RTL direction
            pPr.set('rtl', '1')
            
            # Set right alignment for RTL text
            pPr.set('algn', 'r')
            
            self.stats["rtl_paragraphs_set"] += 1
            
        except Exception as e:
            # Silently fail if RTL setting doesn't work
            pass
    
    def enable_auto_shrink(self, text_frame):
        """
        Enable auto-shrink to prevent text overflow.
        Text will automatically shrink to fit the shape.
        
        Args:
            text_frame: PowerPoint text frame object
        """
        try:
            # Enable word wrap
            text_frame.word_wrap = True
            
            # Enable auto-shrink (resize text to fit shape)
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            self.stats["auto_shrink_enabled"] += 1
            
        except Exception as e:
            # Silently fail if auto-shrink setting doesn't work
            pass
    
    def update_text_runs(self, paragraph, translated_runs: list):
        """
        Update text runs in a paragraph with translated text.
        Preserves all formatting.
        
        Args:
            paragraph: PowerPoint paragraph object
            translated_runs: List of translated run dictionaries
        """
        # Get existing runs
        existing_runs = list(paragraph.runs)
        
        # If counts match, update in place (safest)
        if len(existing_runs) == len(translated_runs):
            for idx, (ppt_run, json_run) in enumerate(zip(existing_runs, translated_runs)):
                ppt_run.text = json_run.get("text", "")
                self.stats["text_runs_updated"] += 1
        else:
            # Counts don't match - need to rebuild runs
            # Clear all runs except first
            while len(paragraph.runs) > 1:
                paragraph._element.remove(paragraph.runs[-1]._element)
            
            # Update first run or create if none
            if len(paragraph.runs) == 0:
                paragraph.add_run()
            
            # Set first run
            if len(translated_runs) > 0:
                paragraph.runs[0].text = translated_runs[0].get("text", "")
                self.stats["text_runs_updated"] += 1
            
            # Add remaining runs
            for json_run in translated_runs[1:]:
                new_run = paragraph.add_run()
                new_run.text = json_run.get("text", "")
                
                # Try to apply formatting from JSON
                try:
                    font = new_run.font
                    if json_run.get("bold") is not None:
                        font.bold = json_run["bold"]
                    if json_run.get("italic") is not None:
                        font.italic = json_run["italic"]
                    if json_run.get("font_size"):
                        font.size = Pt(json_run["font_size"])
                except:
                    pass
                
                self.stats["text_runs_updated"] += 1
        
        # Apply RTL if needed
        self.set_rtl_if_needed(paragraph)
    
    def update_text_frame(self, shape, translated_element: dict):
        """
        Update text in a text frame (TextBox, AutoShape, etc.).
        
        Args:
            shape: PowerPoint shape with text_frame
            translated_element: Translated element dictionary from JSON
        """
        if not shape.has_text_frame:
            return
        
        text_frame = shape.text_frame
        translated_paragraphs = translated_element.get("paragraphs", [])
        
        # Get existing paragraphs
        existing_paragraphs = list(text_frame.paragraphs)
        
        # Update existing paragraphs
        for idx, translated_para in enumerate(translated_paragraphs):
            if idx < len(existing_paragraphs):
                # Update existing paragraph
                ppt_para = existing_paragraphs[idx]
                translated_runs = translated_para.get("runs", [])
                self.update_text_runs(ppt_para, translated_runs)
            else:
                # Add new paragraph
                ppt_para = text_frame.add_paragraph()
                translated_runs = translated_para.get("runs", [])
                self.update_text_runs(ppt_para, translated_runs)
        
        # Remove extra paragraphs if JSON has fewer
        while len(text_frame.paragraphs) > len(translated_paragraphs):
            # Remove last paragraph
            try:
                text_frame._element.remove(text_frame.paragraphs[-1]._element)
            except:
                break
        
        # Enable auto-shrink to prevent overflow
        self.enable_auto_shrink(text_frame)
    
    def update_table(self, shape, translated_table_data: dict):
        """
        Update table cell text with translations.
        
        Args:
            shape: PowerPoint shape with table
            translated_table_data: Translated table_data dictionary from JSON
        """
        if not shape.has_table:
            return
        
        table = shape.table
        translated_cells = translated_table_data.get("cells", [])
        
        # Update each cell
        for cell_data in translated_cells:
            row = cell_data.get("row")
            col = cell_data.get("column")
            
            if row is not None and col is not None:
                try:
                    cell = table.cell(row, col)
                    translated_paragraphs = cell_data.get("paragraphs", [])
                    
                    # Update cell text frame
                    if cell.text_frame and translated_paragraphs:
                        for idx, translated_para in enumerate(translated_paragraphs):
                            if idx < len(cell.text_frame.paragraphs):
                                ppt_para = cell.text_frame.paragraphs[idx]
                                translated_runs = translated_para.get("runs", [])
                                self.update_text_runs(ppt_para, translated_runs)
                        
                        # Enable auto-shrink for table cells
                        self.enable_auto_shrink(cell.text_frame)
                
                except Exception as e:
                    print(f"  Warning: Could not update cell ({row}, {col}): {e}")
        
        self.stats["tables_updated"] += 1
    
    def update_chart(self, shape, translated_chart_data: dict):
        """
        Update chart text elements (title, series names, categories, axis titles, data labels).
        Note: Chart data values are NOT translated, only text labels.
        
        Args:
            shape: PowerPoint shape with chart
            translated_chart_data: Translated chart_data dictionary from JSON
        """
        if not shape.has_chart:
            return
        
        try:
            chart = shape.chart
            
            # Update chart title
            if translated_chart_data.get("title") and chart.has_title:
                try:
                    chart.chart_title.text_frame.text = translated_chart_data["title"]
                    # Enable auto-shrink for chart title
                    self.enable_auto_shrink(chart.chart_title.text_frame)
                    # Apply RTL if needed
                    if self.is_rtl and chart.chart_title.text_frame.paragraphs:
                        for para in chart.chart_title.text_frame.paragraphs:
                            self.set_rtl_if_needed(para)
                except:
                    pass
            
            # Update axis titles
            translated_axis_titles = translated_chart_data.get("axis_titles", {})
            if translated_axis_titles:
                try:
                    # Category axis (X-axis)
                    if "category" in translated_axis_titles and hasattr(chart, 'category_axis'):
                        if chart.category_axis and hasattr(chart.category_axis, 'has_title'):
                            if chart.category_axis.has_title:
                                try:
                                    chart.category_axis.axis_title.text_frame.text = translated_axis_titles["category"]
                                    if self.is_rtl and chart.category_axis.axis_title.text_frame.paragraphs:
                                        for para in chart.category_axis.axis_title.text_frame.paragraphs:
                                            self.set_rtl_if_needed(para)
                                except:
                                    pass
                    
                    # Value axis (Y-axis)
                    if "value" in translated_axis_titles and hasattr(chart, 'value_axis'):
                        if chart.value_axis and hasattr(chart.value_axis, 'has_title'):
                            if chart.value_axis.has_title:
                                try:
                                    chart.value_axis.axis_title.text_frame.text = translated_axis_titles["value"]
                                    if self.is_rtl and chart.value_axis.axis_title.text_frame.paragraphs:
                                        for para in chart.value_axis.axis_title.text_frame.paragraphs:
                                            self.set_rtl_if_needed(para)
                                except:
                                    pass
                    
                    # Series axis (for 3D charts)
                    if "series" in translated_axis_titles and hasattr(chart, 'series_axis'):
                        if chart.series_axis and hasattr(chart.series_axis, 'has_title'):
                            if chart.series_axis.has_title:
                                try:
                                    chart.series_axis.axis_title.text_frame.text = translated_axis_titles["series"]
                                    if self.is_rtl and chart.series_axis.axis_title.text_frame.paragraphs:
                                        for para in chart.series_axis.axis_title.text_frame.paragraphs:
                                            self.set_rtl_if_needed(para)
                                except:
                                    pass
                except:
                    pass
            
            # Update series names and data labels
            translated_data_values = translated_chart_data.get("data_values", [])
            if translated_data_values:
                for idx, series in enumerate(chart.series):
                    if idx < len(translated_data_values):
                        try:
                            # Update series name
                            series_name = translated_data_values[idx].get("series_name")
                            if series_name:
                                series.name = series_name
                            
                            # Update data labels
                            data_labels = translated_data_values[idx].get("data_labels", [])
                            if data_labels:
                                for label_info in data_labels:
                                    try:
                                        point_idx = label_info.get("point_index")
                                        label_text = label_info.get("text")
                                        if point_idx is not None and label_text:
                                            point = series.points[point_idx]
                                            if hasattr(point, 'data_label') and point.data_label:
                                                if hasattr(point.data_label, 'text_frame'):
                                                    point.data_label.text_frame.text = label_text
                                                    if self.is_rtl and point.data_label.text_frame.paragraphs:
                                                        for para in point.data_label.text_frame.paragraphs:
                                                            self.set_rtl_if_needed(para)
                                    except:
                                        continue
                        except:
                            pass
            
            # Fallback: Update series names from series_names list
            translated_series_names = translated_chart_data.get("series_names", [])
            if translated_series_names and not translated_data_values:
                for idx, series in enumerate(chart.series):
                    if idx < len(translated_series_names):
                        try:
                            series.name = translated_series_names[idx]
                        except:
                            pass
            
            # Update categories (if they are text)
            translated_categories = translated_chart_data.get("categories", [])
            if translated_categories:
                try:
                    # Categories are usually in the chart data, hard to update directly
                    # Most chart types don't allow direct category updates after creation
                    pass
                except:
                    pass
            
            self.stats["charts_updated"] += 1
            
        except Exception as e:
            print(f"  Warning: Could not update chart: {e}")
    
    def update_speaker_notes(self, slide, translated_notes: dict):
        """
        Update speaker notes for a slide.
        
        Args:
            slide: PowerPoint slide object
            translated_notes: Translated speaker_notes dictionary from JSON
        """
        if not translated_notes or not translated_notes.get("text"):
            return
        
        try:
            # Ensure slide has notes
            if not slide.has_notes_slide:
                notes_slide = slide.notes_slide  # This creates it
            else:
                notes_slide = slide.notes_slide
            
            # Update notes text
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = translated_notes["text"]
            
            # Apply RTL to notes if needed
            if self.is_rtl and notes_text_frame.paragraphs:
                for para in notes_text_frame.paragraphs:
                    self.set_rtl_if_needed(para)
            
            # Enable auto-shrink for notes
            self.enable_auto_shrink(notes_text_frame)
            
            self.stats["notes_updated"] += 1
            
        except Exception as e:
            print(f"  Warning: Could not update speaker notes: {e}")
    
    def update_slide(self, slide_idx: int):
        """
        Update a single slide with translated content.
        
        Args:
            slide_idx: 0-based slide index
        """
        # Get PowerPoint slide (0-based)
        ppt_slide = self.presentation.slides[slide_idx]
        
        # Get translated slide data (slides array is 0-based in JSON)
        translated_slide = self.translated_data["slides"][slide_idx]
        
        print(f"Processing slide {slide_idx + 1}/{len(self.presentation.slides)}...", end=" ")
        
        # Mirror slide layout for RTL languages (BEFORE updating text)
        self.mirror_slide_layout(ppt_slide)
        
        # Update elements
        translated_elements = translated_slide.get("elements", [])
        
        for element in translated_elements:
            shape_id = element.get("shape_id")
            element_type = element.get("element_type")
            
            # Find shape by ID
            shape = self.find_shape_by_id(ppt_slide, shape_id)
            
            if not shape:
                continue
            
            # Update based on element type
            try:
                if element_type in ["TextBox", "AutoShape"]:
                    self.update_text_frame(shape, element)
                    self.stats["elements_updated"] += 1
                
                elif element_type == "Table":
                    table_data = element.get("table_data")
                    if table_data:
                        self.update_table(shape, table_data)
                        self.stats["elements_updated"] += 1
                
                elif element_type == "Chart":
                    chart_data = element.get("chart_data")
                    if chart_data:
                        self.update_chart(shape, chart_data)
                        self.stats["elements_updated"] += 1
                
                # Picture and other types don't need text updates
                
            except Exception as e:
                print(f"\n  Warning: Error updating shape {shape_id} ({element_type}): {e}")
        
        # Update speaker notes
        translated_notes = translated_slide.get("speaker_notes")
        if translated_notes:
            self.update_speaker_notes(ppt_slide, translated_notes)
        
        self.stats["slides_processed"] += 1
        print("✓")
    
    def verify_slide_count(self):
        """
        Verify that original PPT and translated JSON have same number of slides.
        
        Returns:
            bool: True if counts match
        """
        ppt_count = len(self.presentation.slides)
        json_count = len(self.translated_data["slides"])
        
        if ppt_count != json_count:
            print(f"\n⚠️  WARNING: Slide count mismatch!")
            print(f"   Original PPT: {ppt_count} slides")
            print(f"   Translated JSON: {json_count} slides")
            print(f"   Will process minimum: {min(ppt_count, json_count)} slides")
            return False
        
        return True
    
    def reassemble(self, output_path: str):
        """
        Main reassembly process.
        
        Args:
            output_path: Path to save the reassembled PowerPoint
        """
        print("\n" + "=" * 80)
        print("STARTING REASSEMBLY")
        print("=" * 80)
        print(f"Target language: {self.target_language}")
        if self.is_rtl:
            print(f"RTL mode: ENABLED (Right-to-Left text direction)")
        print(f"Auto-shrink: ENABLED (Text will auto-fit to prevent overflow)")
        
        # Verify slide counts
        self.verify_slide_count()
        
        # Process each slide
        num_slides = min(len(self.presentation.slides), len(self.translated_data["slides"]))
        
        print(f"\nProcessing {num_slides} slides...")
        print("-" * 80)
        
        for slide_idx in range(num_slides):
            self.update_slide(slide_idx)
        
        # Save reassembled presentation
        print("-" * 80)
        print(f"\nSaving reassembled presentation to: {output_path}")
        self.presentation.save(output_path)
        
        # Print statistics
        print("\n" + "=" * 80)
        print("REASSEMBLY COMPLETE!")
        print("=" * 80)
        print(f"Slides processed: {self.stats['slides_processed']}")
        print(f"Elements updated: {self.stats['elements_updated']}")
        print(f"Text runs updated: {self.stats['text_runs_updated']}")
        print(f"Tables updated: {self.stats['tables_updated']}")
        print(f"Charts updated: {self.stats['charts_updated']}")
        print(f"Speaker notes updated: {self.stats['notes_updated']}")
        if self.is_rtl:
            print(f"RTL paragraphs set: {self.stats['rtl_paragraphs_set']}")
            print(f"Shapes mirrored (layout flipped): {self.stats['shapes_mirrored']}")
        print(f"Auto-shrink enabled: {self.stats['auto_shrink_enabled']} text frames")
        print(f"\n✓ Output saved to: {output_path}")
        print("=" * 80)
        
        return self.stats


def main():
    """Main function to run reassembly"""
    parser = argparse.ArgumentParser(
        description="Reassemble PowerPoint from translated JSON",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Basic usage
  python reassembler.py original.pptx translated.json output.pptx
  
  # Using default output name
  python reassembler.py original.pptx translated.json
  
  # The output will be: translated_reassembled.pptx
  
Features:
  - Automatic RTL (Right-to-Left) support for Arabic, Hebrew, etc.
  - Automatic layout mirroring for RTL languages (shapes flip horizontally)
  - Auto-shrink text to prevent overflow
  - Preserves all formatting and visual properties
  - Spanish, French, German, etc. use normal LTR layout
        """
    )
    
    parser.add_argument(
        "original_pptx",
        help="Path to original PowerPoint file (template)"
    )
    parser.add_argument(
        "translated_json",
        help="Path to translated JSON file"
    )
    parser.add_argument(
        "output_pptx",
        nargs='?',
        help="Path to output PowerPoint file (optional, default: <translated_json>_reassembled.pptx)"
    )
    
    args = parser.parse_args()
    
    # Determine output path
    if args.output_pptx:
        output_path = args.output_pptx
    else:
        # Generate default output name
        json_base = args.translated_json.replace(".json", "")
        output_path = f"{json_base}_reassembled.pptx"
    
    # Check if files exist
    if not os.path.exists(args.original_pptx):
        print(f"Error: Original PowerPoint file not found: {args.original_pptx}")
        return 1
    
    if not os.path.exists(args.translated_json):
        print(f"Error: Translated JSON file not found: {args.translated_json}")
        return 1
    
    # Create reassembler and run
    try:
        reassembler = PPTXReassembler(args.original_pptx, args.translated_json)
        stats = reassembler.reassemble(output_path)
        return 0
    except Exception as e:
        print(f"\n❌ Error during reassembly: {e}")
        import traceback
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    exit(main())