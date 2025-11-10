import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE
import zipfile
from lxml import etree
import os

class PPTXExtractor:
    def __init__(self, pptx_path):
        self.pptx_path = pptx_path
        self.presentation = Presentation(pptx_path)
        self.data = {
            "presentation_name": os.path.basename(pptx_path),
            "total_slides": len(self.presentation.slides),
            "slide_masters": [],  # NEW: Store slide master information
            "slides": []
        }
        self.namespaces = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
            'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram',
            'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        }
        
    def get_color_value(self, color_obj):
        """Extract color in multiple formats"""
        color_info = {}
        try:
            if hasattr(color_obj, 'rgb') and color_obj.rgb:
                color_info['rgb'] = str(color_obj.rgb)
            if hasattr(color_obj, 'theme_color') and color_obj.theme_color:
                color_info['theme_color'] = str(color_obj.theme_color)
            if hasattr(color_obj, 'brightness'):
                color_info['brightness'] = color_obj.brightness
        except:
            pass
        return color_info if color_info else None
    
    def extract_background_info(self, slide_or_layout):
        """Extract background information from a slide or layout"""
        background_info = {
            "follows_master": None,
            "fill_type": None,
            "solid_color": None,
            "gradient_colors": None,
            "pattern_type": None,
            "picture_present": False
        }
        
        try:
            # Check if it follows master background
            if hasattr(slide_or_layout, 'follow_master_background'):
                background_info["follows_master"] = slide_or_layout.follow_master_background
            
            # Get background fill
            bg = slide_or_layout.background
            fill = bg.fill
            
            # Determine fill type
            if hasattr(fill, 'type'):
                background_info["fill_type"] = str(fill.type)
            
            # Extract solid fill color
            if hasattr(fill, 'fore_color'):
                try:
                    background_info["solid_color"] = self.get_color_value(fill.fore_color)
                except:
                    pass
            
            # Check for pattern
            if hasattr(fill, 'pattern'):
                try:
                    background_info["pattern_type"] = str(fill.pattern)
                except:
                    pass
            
        except Exception as e:
            # Background might not be accessible
            pass
        
        return background_info
    
    def extract_slide_masters(self):
        """Extract information about all slide masters and their layouts"""
        masters_info = []
        
        for master_idx, master in enumerate(self.presentation.slide_masters):
            master_data = {
                "master_index": master_idx,
                "master_name": getattr(master, 'name', f"Master_{master_idx}"),
                "background": self.extract_background_info(master),
                "layouts": []
            }
            
            # Extract all layouts for this master
            for layout_idx, layout in enumerate(master.slide_layouts):
                layout_data = {
                    "layout_index": layout_idx,
                    "layout_name": layout.name,
                    "background": self.extract_background_info(layout),
                    "placeholders": []
                }
                
                # Extract placeholder information from layout
                try:
                    for placeholder in layout.placeholders:
                        ph_data = {
                            "placeholder_idx": placeholder.placeholder_format.idx,
                            "placeholder_type": str(placeholder.placeholder_format.type),
                            "name": placeholder.name,
                            "dimensions": {
                                "left": placeholder.left,
                                "top": placeholder.top,
                                "width": placeholder.width,
                                "height": placeholder.height
                            }
                        }
                        layout_data["placeholders"].append(ph_data)
                except Exception as e:
                    pass
                
                master_data["layouts"].append(layout_data)
            
            masters_info.append(master_data)
        
        return masters_info
    
    def get_slide_layout_info(self, slide):
        """Get layout information for a specific slide"""
        layout_info = {
            "master_index": None,
            "layout_index": None,
            "layout_name": None,
            "follows_master_background": None
        }
        
        try:
            # Get the slide's layout
            slide_layout = slide.slide_layout
            layout_info["layout_name"] = slide_layout.name
            
            # Check if follows master background
            if hasattr(slide, 'follow_master_background'):
                layout_info["follows_master_background"] = slide.follow_master_background
            
            # Find which master and layout index this corresponds to
            for master_idx, master in enumerate(self.presentation.slide_masters):
                for layout_idx, layout in enumerate(master.slide_layouts):
                    if layout == slide_layout:
                        layout_info["master_index"] = master_idx
                        layout_info["layout_index"] = layout_idx
                        break
                if layout_info["master_index"] is not None:
                    break
                    
        except Exception as e:
            pass
        
        return layout_info
    
    def extract_placeholder_info(self, shape):
        """Extract placeholder-specific information if shape is a placeholder"""
        placeholder_info = {
            "is_placeholder": False,
            "placeholder_type": None,
            "placeholder_idx": None
        }
        
        try:
            if shape.is_placeholder:
                placeholder_info["is_placeholder"] = True
                placeholder_info["placeholder_type"] = str(shape.placeholder_format.type)
                placeholder_info["placeholder_idx"] = shape.placeholder_format.idx
        except:
            pass
        
        return placeholder_info
    
    def extract_shape_fill(self, shape):
        """Extract comprehensive fill information from shape"""
        fill_info = {
            "fill_type": None,
            "solid_color": None,
            "gradient_stops": None,
            "pattern_type": None,
            "picture_present": False
        }
        
        try:
            if hasattr(shape, 'fill'):
                fill = shape.fill
                
                # Fill type
                if hasattr(fill, 'type'):
                    fill_info["fill_type"] = str(fill.type)
                
                # Solid fill
                try:
                    if hasattr(fill, 'fore_color'):
                        fill_info["solid_color"] = self.get_color_value(fill.fore_color)
                except:
                    pass
                
                # Pattern fill
                try:
                    if hasattr(fill, 'pattern'):
                        fill_info["pattern_type"] = str(fill.pattern)
                        if hasattr(fill, 'back_color'):
                            fill_info["pattern_back_color"] = self.get_color_value(fill.back_color)
                except:
                    pass
                
                # Check for gradient
                try:
                    if hasattr(fill, 'gradient_stops'):
                        stops = []
                        for stop in fill.gradient_stops:
                            stops.append({
                                "position": stop.position,
                                "color": self.get_color_value(stop.color)
                            })
                        fill_info["gradient_stops"] = stops
                except:
                    pass
        except:
            pass
        
        return fill_info
    
    def extract_shape_line(self, shape):
        """Extract line/border information from shape"""
        line_info = {
            "has_line": False,
            "color": None,
            "width": None,
            "dash_style": None,
            "transparency": None
        }
        
        try:
            if hasattr(shape, 'line'):
                line = shape.line
                line_info["has_line"] = True
                
                # Line color
                if hasattr(line, 'color'):
                    line_info["color"] = self.get_color_value(line.color)
                
                # Line width
                if hasattr(line, 'width'):
                    line_info["width"] = line.width
                
                # Dash style
                if hasattr(line, 'dash_style'):
                    line_info["dash_style"] = str(line.dash_style)
        except:
            pass
        
        return line_info
    
    def extract_shape_shadow(self, shape):
        """Extract shadow information from shape"""
        shadow_info = {
            "has_shadow": False,
            "shadow_type": None,
            "color": None,
            "transparency": None,
            "blur": None,
            "angle": None,
            "distance": None
        }
        
        try:
            if hasattr(shape, 'shadow'):
                shadow = shape.shadow
                if hasattr(shadow, 'inherit'):
                    shadow_info["has_shadow"] = not shadow.inherit
                
                # Shadow type
                if hasattr(shadow, 'shadow_type'):
                    shadow_info["shadow_type"] = str(shadow.shadow_type)
                
                # Shadow color
                try:
                    if hasattr(shadow, 'color'):
                        shadow_info["color"] = self.get_color_value(shadow.color)
                except:
                    pass
                
                # Shadow properties
                if hasattr(shadow, 'transparency'):
                    shadow_info["transparency"] = shadow.transparency
                if hasattr(shadow, 'blur_radius'):
                    shadow_info["blur"] = shadow.blur_radius
                if hasattr(shadow, 'angle'):
                    shadow_info["angle"] = shadow.angle
                if hasattr(shadow, 'distance'):
                    shadow_info["distance"] = shadow.distance
        except:
            pass
        
        return shadow_info
    
    def extract_bullet_formatting(self, paragraph):
        """Extract comprehensive bullet/numbering information via XML"""
        bullet_info = {
            "is_bulleted": False,
            "bullet_type": None,  # bullet, numbered, none
            "bullet_char": None,
            "bullet_font": None,
            "bullet_color": None,
            "numbering_format": None,  # decimal, alpha, roman
            "start_at": None
        }
        
        try:
            pPr = paragraph._element.pPr
            if pPr is None:
                return bullet_info
            
            # Check for bullet character
            buChar = pPr.find('.//a:buChar', self.namespaces)
            if buChar is not None:
                bullet_info["is_bulleted"] = True
                bullet_info["bullet_type"] = "bullet"
                bullet_info["bullet_char"] = buChar.get('char', '•')
            
            # Check for bullet font
            buFont = pPr.find('.//a:buFont', self.namespaces)
            if buFont is not None:
                bullet_info["bullet_font"] = buFont.get('typeface')
            
            # Check for bullet color
            buClr = pPr.find('.//a:buClr', self.namespaces)
            if buClr is not None:
                # Try to extract color from child elements
                srgbClr = buClr.find('.//a:srgbClr', self.namespaces)
                if srgbClr is not None:
                    bullet_info["bullet_color"] = srgbClr.get('val')
                schemeClr = buClr.find('.//a:schemeClr', self.namespaces)
                if schemeClr is not None:
                    bullet_info["bullet_color"] = f"scheme_{schemeClr.get('val')}"
            
            # Check for numbering
            buAutoNum = pPr.find('.//a:buAutoNum', self.namespaces)
            if buAutoNum is not None:
                bullet_info["is_bulleted"] = True
                bullet_info["bullet_type"] = "numbered"
                num_type = buAutoNum.get('type', 'arabicPeriod')
                bullet_info["numbering_format"] = num_type
                start_at = buAutoNum.get('startAt')
                if start_at:
                    bullet_info["start_at"] = int(start_at)
            
            # Check if bullets are explicitly turned off
            buNone = pPr.find('.//a:buNone', self.namespaces)
            if buNone is not None:
                bullet_info["bullet_type"] = "none"
                bullet_info["is_bulleted"] = False
                
        except Exception as e:
            pass
        
        return bullet_info
    
    def extract_run_formatting(self, run):
        """Extract comprehensive formatting details from a text run including advanced formatting"""
        font = run.font
        formatting = {
            "text": run.text,
            "font_name": font.name,
            "font_size": font.size.pt if font.size else None,
            "bold": font.bold,
            "italic": font.italic,
            "underline": font.underline,
            "color": self.get_color_value(font.color) if hasattr(font, 'color') else None,
            "strike": None,
            "kerning": None,
            "spacing": None,
            "caps": None,
            # NEW: Advanced text formatting
            "superscript": None,
            "subscript": None,
            "text_highlight": None,
            "text_outline": None
        }
        
        # Additional font properties including advanced formatting
        try:
            if hasattr(font, '_element'):
                rPr = run._element.rPr
                if rPr is not None:
                    # Strike through
                    if rPr.find('.//a:strike', self.namespaces) is not None:
                        strike_elem = rPr.find('.//a:strike', self.namespaces)
                        formatting["strike"] = strike_elem.get('val', 'sngStrike')
                    
                    # Kerning
                    if rPr.find('.//a:kern', self.namespaces) is not None:
                        kern_elem = rPr.find('.//a:kern', self.namespaces)
                        formatting["kerning"] = kern_elem.get('val')
                    
                    # Character spacing
                    if rPr.find('.//a:spc', self.namespaces) is not None:
                        spc_elem = rPr.find('.//a:spc', self.namespaces)
                        formatting["spacing"] = spc_elem.get('val')
                    
                    # Capitalization
                    if rPr.get('cap'):
                        formatting["caps"] = rPr.get('cap')
                    
                    # NEW: Superscript/Subscript (baseline)
                    baseline = rPr.get('baseline')
                    if baseline:
                        baseline_val = int(baseline)
                        if baseline_val > 0:
                            formatting["superscript"] = baseline_val
                        elif baseline_val < 0:
                            formatting["subscript"] = abs(baseline_val)
                    
                    # NEW: Text Highlight Color
                    highlight = rPr.find('.//a:highlight', self.namespaces)
                    if highlight is not None:
                        srgbClr = highlight.find('.//a:srgbClr', self.namespaces)
                        if srgbClr is not None:
                            formatting["text_highlight"] = srgbClr.get('val')
                        schemeClr = highlight.find('.//a:schemeClr', self.namespaces)
                        if schemeClr is not None:
                            formatting["text_highlight"] = f"scheme_{schemeClr.get('val')}"
                    
                    # NEW: Character Outline (text border)
                    ln = rPr.find('.//a:ln', self.namespaces)
                    if ln is not None:
                        outline_info = {
                            "width": ln.get('w'),
                            "color": None
                        }
                        
                        # Extract outline color
                        solidFill = ln.find('.//a:solidFill', self.namespaces)
                        if solidFill is not None:
                            srgbClr = solidFill.find('.//a:srgbClr', self.namespaces)
                            if srgbClr is not None:
                                outline_info["color"] = srgbClr.get('val')
                        
                        formatting["text_outline"] = outline_info
                    
        except Exception as e:
            pass
        
        return formatting
    
    def extract_paragraph_formatting(self, paragraph):
        """Extract paragraph-level formatting"""
        para_format = {
            "level": paragraph.level,
            "alignment": str(paragraph.alignment) if paragraph.alignment else None,
            "line_spacing": paragraph.line_spacing,
            "space_before": paragraph.space_before.pt if paragraph.space_before else None,
            "space_after": paragraph.space_after.pt if paragraph.space_after else None,
            "indent": None,  # First line indent
            "left_indent": None,  # Left margin
            "right_indent": None,  # Right margin
            "bullet_format": self.extract_bullet_formatting(paragraph),
            "text_direction": None  # NEW: Text direction for paragraph
        }
        
        # Extract indentation and text direction via XML (python-pptx doesn't expose indent as properties)
        try:
            pPr = paragraph._element.pPr
            if pPr is not None:
                # First line indent
                if pPr.get('indent'):
                    para_format["indent"] = int(pPr.get('indent'))
                # Left margin
                if pPr.get('marL'):
                    para_format["left_indent"] = int(pPr.get('marL'))
                # Right margin
                if pPr.get('marR'):
                    para_format["right_indent"] = int(pPr.get('marR'))
                # Text direction
                rtl = pPr.get('rtl')
                if rtl:
                    para_format["text_direction"] = "rtl" if rtl == '1' else "ltr"
        except:
            pass
        
        return para_format
    
    def extract_text_frame_properties(self, text_frame):
        """Extract text frame properties including text direction and rotation"""
        properties = {
            "margin_left": text_frame.margin_left,
            "margin_right": text_frame.margin_right,
            "margin_top": text_frame.margin_top,
            "margin_bottom": text_frame.margin_bottom,
            "word_wrap": text_frame.word_wrap,
            "auto_size": str(text_frame.auto_size) if hasattr(text_frame, 'auto_size') else None,
            "vertical_anchor": str(text_frame.vertical_anchor) if hasattr(text_frame, 'vertical_anchor') else None,
            "text_direction": None,  # NEW: Text direction (horizontal, vertical, rotated)
            "rotation_angle": None   # NEW: Rotation angle in degrees
        }
        
        # Extract text direction and rotation via XML
        try:
            bodyPr = text_frame._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
            if bodyPr is not None:
                # Text direction (vert, vert270, horz, wordArtVert, etc.)
                vert = bodyPr.get('vert')
                if vert:
                    properties["text_direction"] = vert
                
                # Rotation angle
                rot = bodyPr.get('rot')
                if rot:
                    # Convert from 60000ths of a degree to degrees
                    properties["rotation_angle"] = int(rot) / 60000.0
        except:
            pass
        
        return properties
    
    def extract_table(self, shape):
        """Extract table structure and content"""
        table = shape.table
        table_data = {
            "rows": len(table.rows),
            "columns": len(table.columns),
            "cells": []
        }
        
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                cell_data = {
                    "row": row_idx,
                    "column": col_idx,
                    "text": cell.text,
                    "paragraphs": []
                }
                
                # Extract paragraph and run formatting from cells
                for paragraph in cell.text_frame.paragraphs:
                    para_data = {
                        "paragraph_formatting": self.extract_paragraph_formatting(paragraph),
                        "runs": []
                    }
                    
                    for run in paragraph.runs:
                        run_data = self.extract_run_formatting(run)
                        para_data["runs"].append(run_data)
                    
                    cell_data["paragraphs"].append(para_data)
                
                table_data["cells"].append(cell_data)
        
        return table_data
    
    def extract_chart(self, shape):
        """Extract chart data including chart type, data values, and styling"""
        chart_data = {
            "chart_type": None,      # Chart type (bar, line, pie, etc.)
            "chart_style": None,     # Chart style ID
            "has_title": False,
            "title": None,
            "data_values": [],       # Actual chart data
            "categories": [],
            "series_names": [],
            "axis_titles": {},       # Axis titles (category, value, series)
            "data_labels": [],       # Data labels if present
            "legend_entries": []     # Legend text entries
        }
        
        try:
            chart = shape.chart
            
            # Extract chart type
            if hasattr(chart, 'chart_type'):
                chart_data["chart_type"] = f"{str(chart.chart_type).split('.')[-1]} ({chart.chart_type})"
            
            # Extract chart style
            if hasattr(chart, 'chart_style'):
                chart_data["chart_style"] = chart.chart_style
            
            # Chart title
            if chart.has_title:
                chart_data["has_title"] = True
                try:
                    chart_data["title"] = chart.chart_title.text_frame.text
                except:
                    pass
            
            # Extract axis titles
            try:
                # Category axis (X-axis)
                if hasattr(chart, 'category_axis') and chart.category_axis:
                    if hasattr(chart.category_axis, 'has_title') and chart.category_axis.has_title:
                        try:
                            chart_data["axis_titles"]["category"] = chart.category_axis.axis_title.text_frame.text
                        except:
                            pass
                
                # Value axis (Y-axis)
                if hasattr(chart, 'value_axis') and chart.value_axis:
                    if hasattr(chart.value_axis, 'has_title') and chart.value_axis.has_title:
                        try:
                            chart_data["axis_titles"]["value"] = chart.value_axis.axis_title.text_frame.text
                        except:
                            pass
                
                # Series axis (for 3D charts)
                if hasattr(chart, 'series_axis') and chart.series_axis:
                    if hasattr(chart.series_axis, 'has_title') and chart.series_axis.has_title:
                        try:
                            chart_data["axis_titles"]["series"] = chart.series_axis.axis_title.text_frame.text
                        except:
                            pass
            except Exception as e:
                pass
            
            # Extract data values and data labels from series
            try:
                for series_idx, series in enumerate(chart.series):
                    series_data = {
                        "series_name": series.name,
                        "values": [],
                        "data_labels": []
                    }
                    
                    # Extract values
                    if hasattr(series, 'values'):
                        series_data["values"] = list(series.values)
                    
                    # Extract data labels if present
                    try:
                        if hasattr(series, 'data_labels') and series.data_labels:
                            for point_idx in range(len(series.values) if hasattr(series, 'values') else 0):
                                try:
                                    point = series.points[point_idx]
                                    if hasattr(point, 'data_label') and point.data_label:
                                        if hasattr(point.data_label, 'text_frame'):
                                            label_text = point.data_label.text_frame.text
                                            if label_text:
                                                series_data["data_labels"].append({
                                                    "point_index": point_idx,
                                                    "text": label_text
                                                })
                                except:
                                    continue
                    except:
                        pass
                    
                    chart_data["data_values"].append(series_data)
                    chart_data["series_names"].append(series.name)
            except Exception as e:
                pass
            
            # Extract categories
            try:
                if hasattr(chart, 'plots') and len(chart.plots) > 0:
                    plot = chart.plots[0]
                    if hasattr(plot, 'categories'):
                        chart_data["categories"] = list(plot.categories)
            except:
                pass
            
            # Extract legend entries
            try:
                if hasattr(chart, 'has_legend') and chart.has_legend:
                    legend = chart.legend
                    if hasattr(legend, 'entries'):
                        for entry in legend.entries:
                            try:
                                if hasattr(entry, 'text_frame') and entry.text_frame:
                                    legend_text = entry.text_frame.text
                                    if legend_text:
                                        chart_data["legend_entries"].append(legend_text)
                            except:
                                continue
            except:
                pass
                
        except Exception as e:
            pass
        
        return chart_data
    
    def extract_shape(self, shape, slide_num):
        """Extract comprehensive shape information"""
        element = {
            "shape_id": shape.shape_id,
            "shape_name": shape.name,
            "element_type": None,
            "placeholder_info": self.extract_placeholder_info(shape),  # Placeholder info
            "fill": self.extract_shape_fill(shape),  # NEW: Fill details
            "line": self.extract_shape_line(shape),  # NEW: Border details
            "shadow": self.extract_shape_shadow(shape)  # NEW: Shadow details
        }
        
        # Determine element type
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or (shape.has_text_frame and not shape.has_table):
            element["element_type"] = "TextBox"
            element["text_frame_properties"] = self.extract_text_frame_properties(shape.text_frame)
            element["paragraphs"] = []
            
            for paragraph in shape.text_frame.paragraphs:
                para_data = {
                    "paragraph_formatting": self.extract_paragraph_formatting(paragraph),
                    "runs": []
                }
                
                for run in paragraph.runs:
                    run_data = self.extract_run_formatting(run)
                    para_data["runs"].append(run_data)
                
                element["paragraphs"].append(para_data)
            
            element["full_text"] = shape.text_frame.text
            
        elif shape.has_table:
            element["element_type"] = "Table"
            element["table_data"] = self.extract_table(shape)
            
        elif shape.has_chart:
            element["element_type"] = "Chart"
            element["chart_data"] = self.extract_chart(shape)
        
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            element["element_type"] = "Picture"
            element["image_info"] = {
                "description": shape.name,
                "alt_text": getattr(shape, 'alt_text', None) if hasattr(shape, 'alt_text') else None
            }
        
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            element["element_type"] = "AutoShape"
            if shape.has_text_frame:
                element["text_frame_properties"] = self.extract_text_frame_properties(shape.text_frame)
                element["paragraphs"] = []
                
                for paragraph in shape.text_frame.paragraphs:
                    para_data = {
                        "paragraph_formatting": self.extract_paragraph_formatting(paragraph),
                        "runs": []
                    }
                    
                    for run in paragraph.runs:
                        run_data = self.extract_run_formatting(run)
                        para_data["runs"].append(run_data)
                    
                    element["paragraphs"].append(para_data)
                
                element["full_text"] = shape.text_frame.text
        
        else:
            element["element_type"] = f"Other_{shape.shape_type}"
            # Extract text from any shape that has a text frame, even if it's an "Other" type
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                try:
                    element["text_frame_properties"] = self.extract_text_frame_properties(shape.text_frame)
                    element["paragraphs"] = []
                    
                    for paragraph in shape.text_frame.paragraphs:
                        para_data = {
                            "paragraph_formatting": self.extract_paragraph_formatting(paragraph),
                            "runs": []
                        }
                        
                        for run in paragraph.runs:
                            run_data = self.extract_run_formatting(run)
                            para_data["runs"].append(run_data)
                        
                        element["paragraphs"].append(para_data)
                    
                    element["full_text"] = shape.text_frame.text
                except Exception as e:
                    # If text extraction fails, at least try to get basic text
                    try:
                        if hasattr(shape, 'text'):
                            element["full_text"] = shape.text
                    except:
                        pass
        
        # Dimensions (common for all shapes)
        element["dimensions"] = {
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
            "rotation": shape.rotation
        }
        
        return element
    
    def extract_grouped_shapes(self, group_shape, slide_num):
        """Extract shapes from a group recursively"""
        grouped_elements = []
        
        for shape in group_shape.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                grouped_elements.extend(self.extract_grouped_shapes(shape, slide_num))
            else:
                element = self.extract_shape(shape, slide_num)
                if element:
                    grouped_elements.append(element)
        
        return grouped_elements
    
    def extract_smartart_xml(self):
        """Extract SmartArt diagrams with hierarchical structure using XML parsing"""
        smartart_data = []
        
        try:
            with zipfile.ZipFile(self.pptx_path, 'r') as zip_ref:
                # Find all diagram data files
                diagram_files = [f for f in zip_ref.namelist() if 'diagrams/data' in f and f.endswith('.xml')]
                
                for diagram_file in diagram_files:
                    try:
                        xml_content = zip_ref.read(diagram_file)
                        root = etree.fromstring(xml_content)
                        
                        smartart_element = {
                            "element_type": "SmartArt",
                            "layout_type": None,  # NEW: SmartArt layout type
                            "texts": [],
                            "nodes": [],          # NEW: Hierarchical node structure
                            "full_text": ""
                        }
                        
                        # NEW: Extract layout type from XML
                        try:
                            # Look for layout category in the diagram
                            layout_node = root.find('.//dgm:cat', self.namespaces)
                            if layout_node is not None:
                                smartart_element["layout_type"] = layout_node.get('type')
                            
                            # Alternative: Check prSet for layout info
                            if not smartart_element["layout_type"]:
                                for elem in root.iter():
                                    if 'layoutNode' in elem.tag or 'cat' in elem.tag:
                                        cat_type = elem.get('type')
                                        if cat_type:
                                            smartart_element["layout_type"] = cat_type
                                            break
                        except:
                            pass
                        
                        # NEW: Extract hierarchical node structure
                        try:
                            # Find all points (nodes) in the diagram
                            ptLst = root.find('.//dgm:ptLst', self.namespaces)
                            if ptLst is not None:
                                points = ptLst.findall('.//dgm:pt', self.namespaces)
                                
                                for pt in points:
                                    node_data = {
                                        "node_id": pt.get('modelId'),
                                        "level": None,       # NEW: Hierarchy level
                                        "parent_id": None,   # NEW: Parent node reference
                                        "text": ""
                                    }
                                    
                                    # NEW: Try to determine hierarchy level
                                    prSet = pt.find('.//dgm:prSet', self.namespaces)
                                    if prSet is not None:
                                        # Check for hierarchy level (presLayoutVars or other indicators)
                                        presLayoutVars = prSet.find('.//dgm:presLayoutVars', self.namespaces)
                                        if presLayoutVars is not None:
                                            # Some SmartArt diagrams have depth/level indicators
                                            for child in presLayoutVars:
                                                if 'depth' in child.tag.lower() or 'level' in child.tag.lower():
                                                    try:
                                                        node_data["level"] = int(child.get('val', 0))
                                                    except:
                                                        pass
                                    
                                    # Extract text for this node
                                    t_elem = pt.find('.//dgm:t', self.namespaces)
                                    if t_elem is None:
                                        t_elem = pt.find('.//a:t', self.namespaces)
                                    
                                    if t_elem is not None and t_elem.text:
                                        node_data["text"] = t_elem.text.strip()
                                        smartart_element["texts"].append(node_data["text"])
                                    
                                    smartart_element["nodes"].append(node_data)
                            
                            # NEW: Extract parent-child relationships from cxnLst (connections list)
                            cxnLst = root.find('.//dgm:cxnLst', self.namespaces)
                            if cxnLst is not None:
                                connections = cxnLst.findall('.//dgm:cxn', self.namespaces)
                                
                                for cxn in connections:
                                    cxn_type = cxn.get('type', '')
                                    if cxn_type in ['parOf', 'presOf']:  # Parent-child relationship
                                        src_id = cxn.get('srcId')
                                        dest_id = cxn.get('destId')
                                        
                                        # Update node with parent information
                                        for node in smartart_element["nodes"]:
                                            if node["node_id"] == src_id:
                                                node["parent_id"] = dest_id
                                                break
                            
                            # NEW: If levels weren't found in prSet, infer from parent-child relationships
                            if smartart_element["nodes"]:
                                # Find root nodes (nodes without parents)
                                root_nodes = [n for n in smartart_element["nodes"] if n["parent_id"] is None]
                                
                                # Assign levels based on depth from root
                                def assign_level(node_id, level, nodes):
                                    for node in nodes:
                                        if node["node_id"] == node_id and node["level"] is None:
                                            node["level"] = level
                                            # Find children
                                            children = [n for n in nodes if n["parent_id"] == node_id]
                                            for child in children:
                                                assign_level(child["node_id"], level + 1, nodes)
                                
                                for root_node in root_nodes:
                                    assign_level(root_node["node_id"], 0, smartart_element["nodes"])
                        
                        except Exception as e:
                            # Fallback to simple text extraction if hierarchical extraction fails
                            pass
                        
                        # Fallback: Extract all text if nodes didn't capture everything
                        if not smartart_element["texts"]:
                            for xpath in ['.//dgm:t', './/a:t', './/dgm:text', './/*[local-name()="t"]']:
                                try:
                                    text_elems = root.xpath(xpath, namespaces=self.namespaces)
                                    for elem in text_elems:
                                        if elem.text and elem.text.strip():
                                            smartart_element["texts"].append(elem.text.strip())
                                except:
                                    pass
                        
                        smartart_element["full_text"] = " ".join(smartart_element["texts"])
                        
                        if smartart_element["texts"] or smartart_element["nodes"]:
                            smartart_data.append(smartart_element)
                            
                    except Exception as e:
                        print(f"Error parsing SmartArt file {diagram_file}: {e}")
                        continue
                        
        except Exception as e:
            print(f"SmartArt extraction error: {e}")
        
        return smartart_data
    
    def extract_links(self, shape):
        """Extract hyperlinks from shape"""
        links = []
        
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.hyperlink and run.hyperlink.address:
                        links.append({
                            "text": run.text,
                            "url": run.hyperlink.address
                        })
        
        return links
    
    def extract_slide(self, slide, slide_num):
        """Extract all content from a single slide"""
        slide_data = {
            "slide_number": slide_num,
            "layout_info": self.get_slide_layout_info(slide),  # NEW: Layout information
            "background": self.extract_background_info(slide),  # NEW: Background information
            "elements": [],
            "links": [],
            "speaker_notes": None,
            "smartart": []
        }
        
        # Extract shapes
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                slide_data["elements"].extend(
                    self.extract_grouped_shapes(shape, slide_num)
                )
            else:
                element = self.extract_shape(shape, slide_num)
                if element:
                    slide_data["elements"].append(element)
                
                # Extract links
                links = self.extract_links(shape)
                if links:
                    slide_data["links"].extend(links)
        
        # Extract speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame.text.strip():
                slide_data["speaker_notes"] = {
                    "text": notes_frame.text,
                    "element_type": "SpeakerNotes"
                }
        
        return slide_data
    
    def extract_all(self):
        """Extract all content from presentation"""
        # NEW: Extract slide masters and layouts first
        print("Extracting slide masters and layouts...")
        self.data["slide_masters"] = self.extract_slide_masters()
        
        # Extract slides
        print("Extracting slides...")
        for idx, slide in enumerate(self.presentation.slides, start=1):
            slide_data = self.extract_slide(slide, idx)
            self.data["slides"].append(slide_data)
        
        # Extract SmartArt globally
        print("Extracting SmartArt...")
        smartart_elements = self.extract_smartart_xml()
        
        # Distribute SmartArt to slides (simplified - assign to first slide for now)
        if smartart_elements and self.data["slides"]:
            for smartart in smartart_elements:
                self.data["slides"][0]["smartart"].append(smartart)
        
        return self.data
    
    def save_to_json(self, output_path):
        """Save extracted data to JSON file"""
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(self.data, f, indent=2, ensure_ascii=False)
        print(f"Extraction complete! Saved to {output_path}")


# Usage example
if __name__ == "__main__":
    # Replace with your PPTX file path
    pptx_file = "BI SAM_Negotiations.pptx"
    output_file = "extracted_content_with_layouts.json"
    
    extractor = PPTXExtractor(pptx_file)
    extracted_data = extractor.extract_all()
    extractor.save_to_json(output_file)
    
    print(f"\nTotal slides processed: {extracted_data['total_slides']}")
    print(f"Total slide masters: {len(extracted_data['slide_masters'])}")
    print(f"Total elements extracted: {sum(len(s['elements']) for s in extracted_data['slides'])}")