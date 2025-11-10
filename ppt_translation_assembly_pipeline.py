"""
PowerPoint Translation Assembly Pipeline
==========================================
A comprehensive pipeline that extracts, translates, and reassembles PowerPoint presentations.

This pipeline combines three major components:
1. PPTXExtractor - Extracts all content and metadata from PowerPoint files
2. PPTTranslator - Translates content while preserving 100% of metadata
3. PPTXReassembler - Reassembles translated content back into PowerPoint format

Features:
- Comprehensive metadata preservation
- SmartArt support with hierarchical extraction
- RTL language support with automatic layout mirroring
- Table, Chart, and complex shape handling
- Speaker notes translation
- Background and layout information preservation
"""

import json
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
import zipfile
from lxml import etree
from openai import OpenAI
from typing import Dict, List, Any
from dotenv import load_dotenv
import time
from copy import deepcopy
import argparse


# ============================================================================
# PART 1: PPTX EXTRACTOR
# ============================================================================

class PPTXExtractor:
    def __init__(self, pptx_path):
        self.pptx_path = pptx_path
        self.presentation = Presentation(pptx_path)
        self.data = {
            "presentation_name": os.path.basename(pptx_path),
            "total_slides": len(self.presentation.slides),
            "slide_masters": [],
            "slides": []
        }
        self.namespaces = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
            'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram',
            'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        }
        
    def get_color_value(self, color_obj):
        """Extract color in multiple formats"""
        color_info = {}
        try:
            if hasattr(color_obj, 'rgb') and color_obj.rgb:
                color_info['rgb'] = str(color_obj.rgb)
            if hasattr(color_obj, 'theme_color') and color_obj.theme_color:
                color_info['theme_color'] = str(color_obj.theme_color)
            if hasattr(color_obj, 'brightness'):
                color_info['brightness'] = color_obj.brightness
        except:
            pass
        return color_info if color_info else None
    
    def extract_background_info(self, slide_or_layout):
        """Extract background information from a slide or layout"""
        background_info = {
            "follows_master": None,
            "fill_type": None,
            "solid_color": None,
            "gradient_colors": None,
            "pattern_type": None,
            "picture_present": False
        }
        
        try:
            if hasattr(slide_or_layout, 'follow_master_background'):
                background_info["follows_master"] = slide_or_layout.follow_master_background
            
            bg = slide_or_layout.background
            fill = bg.fill
            
            if hasattr(fill, 'type'):
                background_info["fill_type"] = str(fill.type)
            
            if hasattr(fill, 'fore_color'):
                try:
                    background_info["solid_color"] = self.get_color_value(fill.fore_color)
                except:
                    pass
            
            if hasattr(fill, 'pattern'):
                try:
                    background_info["pattern_type"] = str(fill.pattern)
                except:
                    pass
            
        except Exception as e:
            pass
        
        return background_info
    
    def extract_slide_masters(self):
        """Extract information about all slide masters and their layouts"""
        masters_info = []
        
        for master_idx, master in enumerate(self.presentation.slide_masters):
            master_data = {
                "master_index": master_idx,
                "master_name": getattr(master, 'name', f"Master_{master_idx}"),
                "background": self.extract_background_info(master),
                "layouts": []
            }
            
            for layout_idx, layout in enumerate(master.slide_layouts):
                layout_data = {
                    "layout_index": layout_idx,
                    "layout_name": layout.name,
                    "background": self.extract_background_info(layout),
                    "placeholders": []
                }
                
                try:
                    for placeholder in layout.placeholders:
                        ph_data = {
                            "placeholder_idx": placeholder.placeholder_format.idx,
                            "placeholder_type": str(placeholder.placeholder_format.type),
                            "name": placeholder.name,
                            "dimensions": {
                                "left": placeholder.left,
                                "top": placeholder.top,
                                "width": placeholder.width,
                                "height": placeholder.height
                            }
                        }
                        layout_data["placeholders"].append(ph_data)
                except Exception as e:
                    pass
                
                master_data["layouts"].append(layout_data)
            
            masters_info.append(master_data)
        
        return masters_info
    
    def get_slide_layout_info(self, slide):
        """Get layout information for a specific slide"""
        layout_info = {
            "master_index": None,
            "layout_index": None,
            "layout_name": None,
            "follows_master_background": None
        }
        
        try:
            slide_layout = slide.slide_layout
            layout_info["layout_name"] = slide_layout.name
            
            if hasattr(slide, 'follow_master_background'):
                layout_info["follows_master_background"] = slide.follow_master_background
            
            for master_idx, master in enumerate(self.presentation.slide_masters):
                for layout_idx, layout in enumerate(master.slide_layouts):
                    if layout == slide_layout:
                        layout_info["master_index"] = master_idx
                        layout_info["layout_index"] = layout_idx
                        break
                if layout_info["master_index"] is not None:
                    break
                    
        except Exception as e:
            pass
        
        return layout_info
    
    def extract_placeholder_info(self, shape):
        """Extract placeholder-specific information if shape is a placeholder"""
        placeholder_info = {
            "is_placeholder": False,
            "placeholder_type": None,
            "placeholder_idx": None
        }
        
        try:
            if shape.is_placeholder:
                placeholder_info["is_placeholder"] = True
                placeholder_info["placeholder_type"] = str(shape.placeholder_format.type)
                placeholder_info["placeholder_idx"] = shape.placeholder_format.idx
        except:
            pass
        
        return placeholder_info
    
    def extract_shape_fill(self, shape):
        """Extract comprehensive fill information from shape"""
        fill_info = {
            "fill_type": None,
            "solid_color": None,
            "gradient_stops": None,
            "pattern_type": None,
            "picture_present": False
        }
        
        try:
            if hasattr(shape, 'fill'):
                fill = shape.fill
                
                if hasattr(fill, 'type'):
                    fill_info["fill_type"] = str(fill.type)
                
                try:
                    if hasattr(fill, 'fore_color'):
                        fill_info["solid_color"] = self.get_color_value(fill.fore_color)
                except:
                    pass
                
                try:
                    if hasattr(fill, 'pattern'):
                        fill_info["pattern_type"] = str(fill.pattern)
                        if hasattr(fill, 'back_color'):
                            fill_info["pattern_back_color"] = self.get_color_value(fill.back_color)
                except:
                    pass
                
                try:
                    if hasattr(fill, 'gradient_stops'):
                        stops = []
                        for stop in fill.gradient_stops:
                            stops.append({
                                "position": stop.position,
                                "color": self.get_color_value(stop.color)
                            })
                        fill_info["gradient_stops"] = stops
                except:
                    pass
        except:
            pass
        
        return fill_info
    
    def extract_shape_line(self, shape):
        """Extract line/border information from shape"""
        line_info = {
            "has_line": False,
            "color": None,
            "width": None,
            "dash_style": None,
            "transparency": None
        }
        
        try:
            if hasattr(shape, 'line'):
                line = shape.line
                line_info["has_line"] = True
                
                if hasattr(line, 'color'):
                    line_info["color"] = self.get_color_value(line.color)
                
                if hasattr(line, 'width'):
                    line_info["width"] = line.width
                
                if hasattr(line, 'dash_style'):
                    line_info["dash_style"] = str(line.dash_style)
        except:
            pass
        
        return line_info
    
    def extract_shape_shadow(self, shape):
        """Extract shadow information from shape"""
        shadow_info = {
            "has_shadow": False,
            "shadow_type": None,
            "color": None,
            "transparency": None,
            "blur": None,
            "angle": None,
            "distance": None
        }
        
        try:
            if hasattr(shape, 'shadow'):
                shadow = shape.shadow
                if hasattr(shadow, 'inherit'):
                    shadow_info["has_shadow"] = not shadow.inherit
                
                if hasattr(shadow, 'shadow_type'):
                    shadow_info["shadow_type"] = str(shadow.shadow_type)
                
                try:
                    if hasattr(shadow, 'color'):
                        shadow_info["color"] = self.get_color_value(shadow.color)
                except:
                    pass
                
                if hasattr(shadow, 'transparency'):
                    shadow_info["transparency"] = shadow.transparency
                if hasattr(shadow, 'blur_radius'):
                    shadow_info["blur"] = shadow.blur_radius
                if hasattr(shadow, 'angle'):
                    shadow_info["angle"] = shadow.angle
                if hasattr(shadow, 'distance'):
                    shadow_info["distance"] = shadow.distance
        except:
            pass
        
        return shadow_info
    
    def extract_bullet_formatting(self, paragraph):
        """Extract comprehensive bullet/numbering information via XML"""
        bullet_info = {
            "is_bulleted": False,
            "bullet_type": None,
            "bullet_char": None,
            "bullet_font": None,
            "bullet_color": None,
            "numbering_format": None,
            "start_at": None
        }
        
        try:
            pPr = paragraph._element.pPr
            if pPr is None:
                return bullet_info
            
            buChar = pPr.find('.//a:buChar', self.namespaces)
            if buChar is not None:
                bullet_info["is_bulleted"] = True
                bullet_info["bullet_type"] = "bullet"
                bullet_info["bullet_char"] = buChar.get('char', '•')
            
            buFont = pPr.find('.//a:buFont', self.namespaces)
            if buFont is not None:
                bullet_info["bullet_font"] = buFont.get('typeface')
            
            buClr = pPr.find('.//a:buClr', self.namespaces)
            if buClr is not None:
                srgbClr = buClr.find('.//a:srgbClr', self.namespaces)
                if srgbClr is not None:
                    bullet_info["bullet_color"] = srgbClr.get('val')
                schemeClr = buClr.find('.//a:schemeClr', self.namespaces)
                if schemeClr is not None:
                    bullet_info["bullet_color"] = f"scheme_{schemeClr.get('val')}"
            
            buAutoNum = pPr.find('.//a:buAutoNum', self.namespaces)
            if buAutoNum is not None:
                bullet_info["is_bulleted"] = True
                bullet_info["bullet_type"] = "numbered"
                num_type = buAutoNum.get('type', 'arabicPeriod')
                bullet_info["numbering_format"] = num_type
                start_at = buAutoNum.get('startAt')
                if start_at:
                    bullet_info["start_at"] = int(start_at)
            
            buNone = pPr.find('.//a:buNone', self.namespaces)
            if buNone is not None:
                bullet_info["bullet_type"] = "none"
                bullet_info["is_bulleted"] = False
                
        except Exception as e:
            pass
        
        return bullet_info
    
    def extract_run_formatting(self, run):
        """Extract comprehensive formatting details from a text run including advanced formatting"""
        font = run.font
        formatting = {
            "text": run.text,
            "font_name": font.name,
            "font_size": font.size.pt if font.size else None,
            "bold": font.bold,
            "italic": font.italic,
            "underline": font.underline,
            "color": self.get_color_value(font.color) if hasattr(font, 'color') else None,
            "strike": None,
            "kerning": None,
            "spacing": None,
            "caps": None,
            "superscript": None,
            "subscript": None,
            "text_highlight": None,
            "text_outline": None
        }
        
        try:
            if hasattr(font, '_element'):
                rPr = run._element.rPr
                if rPr is not None:
                    if rPr.find('.//a:strike', self.namespaces) is not None:
                        strike_elem = rPr.find('.//a:strike', self.namespaces)
                        formatting["strike"] = strike_elem.get('val', 'sngStrike')
                    
                    if rPr.find('.//a:kern', self.namespaces) is not None:
                        kern_elem = rPr.find('.//a:kern', self.namespaces)
                        formatting["kerning"] = int(kern_elem.get('val', 0))
                    
                    if rPr.find('.//a:spc', self.namespaces) is not None:
                        spc_elem = rPr.find('.//a:spc', self.namespaces)
                        formatting["spacing"] = int(spc_elem.get('val', 0))
                    
                    if rPr.find('.//a:cap', self.namespaces) is not None:
                        cap_elem = rPr.find('.//a:cap', self.namespaces)
                        formatting["caps"] = cap_elem.get('cap', 'none')
                    
                    baseline_elem = rPr.find('.//a:baseline', self.namespaces)
                    if baseline_elem is not None:
                        baseline_val = int(baseline_elem.get('val', 0))
                        if baseline_val > 0:
                            formatting["superscript"] = True
                        elif baseline_val < 0:
                            formatting["subscript"] = True
                    
                    highlight_elem = rPr.find('.//a:highlight', self.namespaces)
                    if highlight_elem is not None:
                        formatting["text_highlight"] = True
                    
                    outline_elem = rPr.find('.//a:ln', self.namespaces)
                    if outline_elem is not None:
                        formatting["text_outline"] = True
        except:
            pass
        
        return formatting
    
    def extract_paragraph_formatting(self, paragraph):
        """Extract paragraph-level formatting"""
        formatting = {
            "alignment": str(paragraph.alignment) if paragraph.alignment else None,
            "level": paragraph.level,
            "line_spacing": paragraph.line_spacing,
            "space_before": paragraph.space_before.pt if paragraph.space_before else None,
            "space_after": paragraph.space_after.pt if paragraph.space_after else None,
            "bullet_info": self.extract_bullet_formatting(paragraph)
        }
        
        return formatting
    
    def extract_text_content(self, shape):
        """Extract text content with full formatting preservation"""
        if not shape.has_text_frame:
            return None
        
        text_frame = shape.text_frame
        paragraphs = []
        
        for para in text_frame.paragraphs:
            para_data = {
                "formatting": self.extract_paragraph_formatting(para),
                "runs": []
            }
            
            for run in para.runs:
                run_data = self.extract_run_formatting(run)
                para_data["runs"].append(run_data)
            
            paragraphs.append(para_data)
        
        return {
            "paragraphs": paragraphs,
            "full_text": text_frame.text
        }
    
    def extract_table(self, shape):
        """Extract table data with cell formatting"""
        if not shape.has_table:
            return None
        
        table = shape.table
        table_data = {
            "rows": len(table.rows),
            "columns": len(table.columns),
            "cells": []
        }
        
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                cell_data = {
                    "row": row_idx,
                    "column": col_idx,
                    "text": cell.text,
                    "paragraphs": []
                }
                
                if cell.text_frame:
                    for para in cell.text_frame.paragraphs:
                        para_data = {
                            "formatting": self.extract_paragraph_formatting(para),
                            "runs": []
                        }
                        
                        for run in para.runs:
                            run_data = self.extract_run_formatting(run)
                            para_data["runs"].append(run_data)
                        
                        cell_data["paragraphs"].append(para_data)
                
                table_data["cells"].append(cell_data)
        
        return table_data
    
    def extract_chart(self, shape):
        """Extract chart data and metadata"""
        if not shape.has_chart:
            return None
        
        chart = shape.chart
        chart_data = {
            "chart_type": str(chart.chart_type),
            "has_title": chart.has_title,
            "title": chart.chart_title.text_frame.text if chart.has_title else None,
            "series": [],
            "categories": []
        }
        
        try:
            for series in chart.series:
                series_data = {
                    "name": series.name,
                    "values": [v for v in series.values] if hasattr(series, 'values') else []
                }
                chart_data["series"].append(series_data)
        except:
            pass
        
        try:
            if hasattr(chart, 'plots') and len(chart.plots) > 0:
                plot = chart.plots[0]
                if hasattr(plot, 'categories'):
                    chart_data["categories"] = [cat for cat in plot.categories]
        except:
            pass
        
        return chart_data
    
    def extract_shape(self, shape, slide_num):
        """Extract all information from a single shape"""
        element = {
            "shape_id": shape.shape_id,
            "shape_name": shape.name,
            "element_type": str(shape.shape_type).split('.')[-1],
            "position": {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            },
            "rotation": shape.rotation,
            "placeholder_info": self.extract_placeholder_info(shape),
            "fill": self.extract_shape_fill(shape),
            "line": self.extract_shape_line(shape),
            "shadow": self.extract_shape_shadow(shape)
        }
        
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            element["element_type"] = "TextBox" if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX else "AutoShape"
            text_content = self.extract_text_content(shape)
            if text_content:
                element.update(text_content)
        
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            element["element_type"] = "Table"
            table_data = self.extract_table(shape)
            if table_data:
                element["table_data"] = table_data
        
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            element["element_type"] = "Chart"
            chart_data = self.extract_chart(shape)
            if chart_data:
                element["chart_data"] = chart_data
        
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            element["element_type"] = "Picture"
            element["image_info"] = {
                "has_image": True,
                "image_type": "embedded"
            }
        
        return element
    
    def extract_grouped_shapes(self, group_shape, slide_num):
        """Extract shapes from a group"""
        grouped_elements = []
        
        for shape in group_shape.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                grouped_elements.extend(
                    self.extract_grouped_shapes(shape, slide_num)
                )
            else:
                element = self.extract_shape(shape, slide_num)
                if element:
                    element["is_grouped"] = True
                    element["group_id"] = group_shape.shape_id
                    grouped_elements.append(element)
        
        return grouped_elements
    
    def extract_smartart_xml(self):
        """Extract SmartArt content from XML with hierarchical node information"""
        smartart_data = []
        
        try:
            with zipfile.ZipFile(self.pptx_path, 'r') as zip_ref:
                diagram_files = [f for f in zip_ref.namelist() if 'diagrams/data' in f and f.endswith('.xml')]
                
                for diagram_file in diagram_files:
                    try:
                        xml_content = zip_ref.read(diagram_file)
                        root = etree.fromstring(xml_content)
                        
                        smartart_element = {
                            "diagram_file": diagram_file,
                            "texts": [],
                            "nodes": [],
                            "full_text": ""
                        }
                        
                        try:
                            ptLst = root.find('.//dgm:ptLst', self.namespaces)
                            if ptLst is not None:
                                points = ptLst.findall('.//dgm:pt', self.namespaces)
                                
                                for pt in points:
                                    node_data = {
                                        "node_id": pt.get('modelId'),
                                        "node_type": pt.get('type'),
                                        "text": None,
                                        "level": None,
                                        "parent_id": None
                                    }
                                    
                                    prSet = pt.find('.//dgm:prSet', self.namespaces)
                                    if prSet is not None:
                                        presLayoutVars = prSet.find('.//dgm:presLayoutVars', self.namespaces)
                                        if presLayoutVars is not None:
                                            for child in presLayoutVars:
                                                if 'depth' in child.tag.lower() or 'level' in child.tag.lower():
                                                    try:
                                                        node_data["level"] = int(child.get('val', 0))
                                                    except:
                                                        pass
                                    
                                    t_elem = pt.find('.//dgm:t', self.namespaces)
                                    if t_elem is None:
                                        t_elem = pt.find('.//a:t', self.namespaces)
                                    
                                    if t_elem is not None and t_elem.text:
                                        node_data["text"] = t_elem.text.strip()
                                        smartart_element["texts"].append(node_data["text"])
                                    
                                    smartart_element["nodes"].append(node_data)
                            
                            cxnLst = root.find('.//dgm:cxnLst', self.namespaces)
                            if cxnLst is not None:
                                connections = cxnLst.findall('.//dgm:cxn', self.namespaces)
                                
                                for cxn in connections:
                                    cxn_type = cxn.get('type', '')
                                    if cxn_type in ['parOf', 'presOf']:
                                        src_id = cxn.get('srcId')
                                        dest_id = cxn.get('destId')
                                        
                                        for node in smartart_element["nodes"]:
                                            if node["node_id"] == src_id:
                                                node["parent_id"] = dest_id
                                                break
                            
                            if smartart_element["nodes"]:
                                root_nodes = [n for n in smartart_element["nodes"] if n["parent_id"] is None]
                                
                                def assign_level(node_id, level, nodes):
                                    for node in nodes:
                                        if node["node_id"] == node_id and node["level"] is None:
                                            node["level"] = level
                                            children = [n for n in nodes if n["parent_id"] == node_id]
                                            for child in children:
                                                assign_level(child["node_id"], level + 1, nodes)
                                
                                for root_node in root_nodes:
                                    assign_level(root_node["node_id"], 0, smartart_element["nodes"])
                        
                        except Exception as e:
                            pass
                        
                        if not smartart_element["texts"]:
                            for xpath in ['.//dgm:t', './/a:t', './/dgm:text', './/*[local-name()="t"]']:
                                try:
                                    text_elems = root.xpath(xpath, namespaces=self.namespaces)
                                    for elem in text_elems:
                                        if elem.text and elem.text.strip():
                                            smartart_element["texts"].append(elem.text.strip())
                                except:
                                    pass
                        
                        smartart_element["full_text"] = " ".join(smartart_element["texts"])
                        
                        if smartart_element["texts"] or smartart_element["nodes"]:
                            smartart_data.append(smartart_element)
                            
                    except Exception as e:
                        print(f"Error parsing SmartArt file {diagram_file}: {e}")
                        continue
                        
        except Exception as e:
            print(f"SmartArt extraction error: {e}")
        
        return smartart_data
    
    def extract_links(self, shape):
        """Extract hyperlinks from shape"""
        links = []
        
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.hyperlink and run.hyperlink.address:
                        links.append({
                            "text": run.text,
                            "url": run.hyperlink.address
                        })
        
        return links
    
    def extract_slide(self, slide, slide_num):
        """Extract all content from a single slide"""
        slide_data = {
            "slide_number": slide_num,
            "layout_info": self.get_slide_layout_info(slide),
            "background": self.extract_background_info(slide),
            "elements": [],
            "links": [],
            "speaker_notes": None,
            "smartart": []
        }
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                slide_data["elements"].extend(
                    self.extract_grouped_shapes(shape, slide_num)
                )
            else:
                element = self.extract_shape(shape, slide_num)
                if element:
                    slide_data["elements"].append(element)
                
                links = self.extract_links(shape)
                if links:
                    slide_data["links"].extend(links)
        
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame.text.strip():
                slide_data["speaker_notes"] = {
                    "text": notes_frame.text,
                    "element_type": "SpeakerNotes"
                }
        
        return slide_data
    
    def extract_all(self):
        """Extract all content from presentation"""
        print("Extracting slide masters and layouts...")
        self.data["slide_masters"] = self.extract_slide_masters()
        
        print("Extracting slides...")
        for idx, slide in enumerate(self.presentation.slides, start=1):
            slide_data = self.extract_slide(slide, idx)
            self.data["slides"].append(slide_data)
        
        print("Extracting SmartArt...")
        smartart_elements = self.extract_smartart_xml()
        
        if smartart_elements and self.data["slides"]:
            for smartart in smartart_elements:
                self.data["slides"][0]["smartart"].append(smartart)
        
        return self.data
    
    def save_to_json(self, output_path):
        """Save extracted data to JSON file"""
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(self.data, f, indent=2, ensure_ascii=False)
        print(f"Extraction complete! Saved to {output_path}")


# ============================================================================
# PART 2: PPT TRANSLATOR
# ============================================================================

class PPTTranslator:
    """
    Translates PowerPoint extracted content while preserving 100% of metadata.
    Only translates actual text content, keeping all formatting and structural data intact.
    """
    
    def __init__(self, api_key: str = None, target_language: str = "Spanish"):
        """
        Initialize the translator.
        
        Args:
            api_key: OpenAI API key (if None, loads from .env)
            target_language: Target language for translation (default: Spanish)
        """
        load_dotenv()
        
        self.api_key = api_key or os.getenv('OPENAI_API_KEY')
        if not self.api_key:
            raise ValueError("OPENAI_API_KEY not found. Please set it in .env file or pass it as parameter.")
        
        self.client = OpenAI(api_key=self.api_key)
        self.target_language = target_language
        
        self.model = "gpt-4o-mini"
        
        # Check if target language is RTL
        rtl_languages = ['arabic', 'hebrew', 'urdu', 'persian', 'farsi', 'pashto', 'sindhi', 'yiddish']
        self.is_rtl = any(lang in target_language.lower() for lang in rtl_languages)
        
        self.stats = {
            "total_texts_translated": 0,
            "api_calls": 0,
            "total_tokens_used": 0
        }
    
    def translate_batch(self, texts: List[str]) -> List[str]:
        """
        Translate a batch of texts using GPT-4o-mini.
        
        Args:
            texts: List of text strings to translate
            
        Returns:
            List of translated text strings in the same order
        """
        if not texts:
            return []
        
        text_map = {}
        non_empty_texts = []
        for idx, text in enumerate(texts):
            if text and text.strip():
                text_map[len(non_empty_texts)] = idx
                non_empty_texts.append(text)
        
        if not non_empty_texts:
            return texts
        
        texts_json = []
        for idx, text in enumerate(non_empty_texts):
            texts_json.append({
                "id": idx,
                "text": text
            })
        
        import json as json_module
        batch_json = json_module.dumps(texts_json, ensure_ascii=False)
        
        prompt = f"""Translate the texts in the following JSON array to {self.target_language}.

CRITICAL RULES:
1. Return ONLY a JSON array with the same structure
2. Keep the same "id" values
3. Translate only the "text" field
4. Preserve all line breaks (\\n) and special characters
5. Do not add any explanations or extra content outside the JSON
6. The number of items in output must match the input exactly

Input JSON:
{batch_json}

Output (JSON array only):"""

        try:
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": f"You are a professional translator. Return only valid JSON. Translate to {self.target_language}."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.3,
                max_tokens=4096
            )
            
            self.stats["api_calls"] += 1
            self.stats["total_tokens_used"] += response.usage.total_tokens
            
            response_text = response.choices[0].message.content.strip()
            
            if "```json" in response_text:
                response_text = response_text.split("```json")[1].split("```")[0].strip()
            elif "```" in response_text:
                response_text = response_text.split("```")[1].split("```")[0].strip()
            
            if not response_text.startswith('['):
                start_idx = response_text.find('[')
                if start_idx != -1:
                    response_text = response_text[start_idx:]
            
            if not response_text.endswith(']'):
                end_idx = response_text.rfind(']')
                if end_idx != -1:
                    response_text = response_text[:end_idx + 1]
            
            try:
                translated_json = json_module.loads(response_text)
            except json_module.JSONDecodeError:
                import re
                match = re.search(r'\[.*\]', response_text, re.DOTALL)
                if match:
                    response_text = match.group(0)
                    translated_json = json_module.loads(response_text)
                else:
                    raise
            
            if not isinstance(translated_json, list) or len(translated_json) != len(non_empty_texts):
                if isinstance(translated_json, list) and all(isinstance(item, dict) and 'id' in item for item in translated_json):
                    translated_json = sorted(translated_json, key=lambda x: x.get('id', 0))
                    translated_texts = [item.get('text', non_empty_texts[i]) for i, item in enumerate(translated_json[:len(non_empty_texts)])]
                else:
                    return texts
            else:
                translated_json = sorted(translated_json, key=lambda x: x.get('id', 0))
                translated_texts = [item.get('text', '') for item in translated_json]
            
            result = texts.copy()
            for new_idx, orig_idx in text_map.items():
                if new_idx < len(translated_texts):
                    result[orig_idx] = translated_texts[new_idx]
            
            self.stats["total_texts_translated"] += len(non_empty_texts)
            return result
            
        except json_module.JSONDecodeError:
            return self.translate_one_by_one(texts)
        except Exception as e:
            print(f"Translation error: {e}")
            return self.translate_one_by_one(texts)
    
    def translate_one_by_one(self, texts: List[str]) -> List[str]:
        """
        Fallback method: translate texts one by one.
        
        Args:
            texts: List of text strings to translate
            
        Returns:
            List of translated text strings
        """
        translated = []
        for text in texts:
            if not text or not text.strip():
                translated.append(text)
                continue
            
            try:
                response = self.client.chat.completions.create(
                    model=self.model,
                    messages=[
                        {"role": "system", "content": f"You are a professional translator. Translate to {self.target_language}. Return ONLY the translated text, nothing else."},
                        {"role": "user", "content": f"Translate this to {self.target_language}:\n\n{text}"}
                    ],
                    temperature=0.3,
                    max_tokens=2048
                )
                
                self.stats["api_calls"] += 1
                self.stats["total_tokens_used"] += response.usage.total_tokens
                self.stats["total_texts_translated"] += 1
                
                translated_text = response.choices[0].message.content.strip()
                translated.append(translated_text)
                
            except Exception as e:
                print(f"Error translating text: {e}")
                translated.append(text)
        
        return translated
    
    def translate_text_element(self, element: Dict) -> Dict:
        """
        Translate a text element (TextBox, AutoShape) while preserving all formatting.
        
        Args:
            element: Element dictionary with paragraphs
            
        Returns:
            Element dictionary with translated text
        """
        new_element = deepcopy(element)
        
        if "paragraphs" not in new_element:
            return new_element
        
        for paragraph in new_element["paragraphs"]:
            if "runs" not in paragraph:
                continue
            
            texts_to_translate = [run.get("text", "") for run in paragraph["runs"]]
            
            translated_texts = self.translate_batch(texts_to_translate)
            
            for run, translated_text in zip(paragraph["runs"], translated_texts):
                run["text"] = translated_text
        
        if "full_text" in new_element:
            all_texts = []
            for para in new_element["paragraphs"]:
                para_texts = [run.get("text", "") for run in para.get("runs", [])]
                all_texts.extend(para_texts)
            new_element["full_text"] = "".join(all_texts)
        
        return new_element
    
    def translate_table(self, table_data: Dict) -> Dict:
        """
        Translate table cell content while preserving structure.
        
        Args:
            table_data: Table data dictionary
            
        Returns:
            Table data dictionary with translated content
        """
        new_table = deepcopy(table_data)
        
        if "cells" not in new_table:
            return new_table
        
        for cell in new_table["cells"]:
            if "paragraphs" not in cell:
                continue
            
            for paragraph in cell["paragraphs"]:
                if "runs" not in paragraph:
                    continue
                
                texts_to_translate = [run.get("text", "") for run in paragraph["runs"]]
                translated_texts = self.translate_batch(texts_to_translate)
                
                for run, translated_text in zip(paragraph["runs"], translated_texts):
                    run["text"] = translated_text
            
            if "text" in cell:
                all_texts = []
                for para in cell["paragraphs"]:
                    para_texts = [run.get("text", "") for run in para.get("runs", [])]
                    all_texts.extend(para_texts)
                cell["text"] = "".join(all_texts)
        
        return new_table
    
    def translate_chart(self, chart_data: Dict) -> Dict:
        """
        Translate chart text elements (title, series names, categories).
        Note: Chart data values are NOT translated.
        
        Args:
            chart_data: Chart data dictionary
            
        Returns:
            Chart data dictionary with translated text
        """
        new_chart = deepcopy(chart_data)
        
        texts_to_translate = []
        
        if new_chart.get("title"):
            texts_to_translate.append(new_chart["title"])
        
        for series in new_chart.get("series", []):
            if series.get("name"):
                texts_to_translate.append(series["name"])
        
        for category in new_chart.get("categories", []):
            if isinstance(category, str):
                texts_to_translate.append(category)
        
        if texts_to_translate:
            translated_texts = self.translate_batch(texts_to_translate)
            
            idx = 0
            if new_chart.get("title"):
                new_chart["title"] = translated_texts[idx]
                idx += 1
            
            for series in new_chart.get("series", []):
                if series.get("name"):
                    series["name"] = translated_texts[idx]
                    idx += 1
            
            translated_categories = []
            for category in new_chart.get("categories", []):
                if isinstance(category, str):
                    translated_categories.append(translated_texts[idx])
                    idx += 1
                else:
                    translated_categories.append(category)
            new_chart["categories"] = translated_categories
        
        return new_chart
    
    def translate_smartart(self, smartart: Dict) -> Dict:
        """
        Translate SmartArt content while preserving hierarchy.
        
        Args:
            smartart: SmartArt dictionary
            
        Returns:
            SmartArt dictionary with translated content
        """
        new_smartart = deepcopy(smartart)
        
        if "texts" in new_smartart and new_smartart["texts"]:
            translated_texts = self.translate_batch(new_smartart["texts"])
            new_smartart["texts"] = translated_texts
        
        if "nodes" in new_smartart and new_smartart["nodes"]:
            node_texts = [node.get("text", "") for node in new_smartart["nodes"]]
            translated_node_texts = self.translate_batch(node_texts)
            
            for idx, node in enumerate(new_smartart["nodes"]):
                if node.get("text"):
                    node["text"] = translated_node_texts[idx]
        
        if "texts" in new_smartart:
            new_smartart["full_text"] = " ".join(new_smartart["texts"])
        
        return new_smartart
    
    def translate_speaker_notes(self, notes: Dict) -> Dict:
        """
        Translate speaker notes while preserving metadata.
        
        Args:
            notes: Speaker notes dictionary
            
        Returns:
            Speaker notes dictionary with translated text
        """
        new_notes = deepcopy(notes)
        
        if "text" in new_notes and new_notes["text"]:
            translated = self.translate_batch([new_notes["text"]])
            new_notes["text"] = translated[0]
        
        return new_notes
    
    def translate_slide(self, slide: Dict, slide_num: int) -> Dict:
        """
        Translate a single slide while preserving all metadata.
        
        Args:
            slide: Slide dictionary
            slide_num: Slide number for progress reporting
            
        Returns:
            Slide dictionary with translated content
        """
        print(f"Translating slide {slide_num}...", end=" ", flush=True)
        
        if slide_num > 1:
            time.sleep(0.2)
        
        new_slide = deepcopy(slide)
        
        if "elements" in new_slide:
            translated_elements = []
            for element in new_slide["elements"]:
                element_type = element.get("element_type")
                
                if element_type == "Table":
                    if "table_data" in element:
                        element["table_data"] = self.translate_table(element["table_data"])
                    translated_elements.append(element)
                    
                elif element_type == "Chart":
                    if "chart_data" in element:
                        element["chart_data"] = self.translate_chart(element["chart_data"])
                    translated_elements.append(element)
                    
                elif element_type in ["TextBox", "AutoShape"]:
                    translated_elements.append(self.translate_text_element(element))
                    
                else:
                    translated_elements.append(deepcopy(element))
            
            new_slide["elements"] = translated_elements
        
        if "speaker_notes" in new_slide and new_slide["speaker_notes"]:
            new_slide["speaker_notes"] = self.translate_speaker_notes(new_slide["speaker_notes"])
        
        if "smartart" in new_slide and new_slide["smartart"]:
            translated_smartart = []
            for smartart in new_slide["smartart"]:
                translated_smartart.append(self.translate_smartart(smartart))
            new_slide["smartart"] = translated_smartart
        
        print("✓")
        return new_slide
    
    def translate_presentation(self, input_path: str, output_path: str) -> Dict:
        """
        Translate entire presentation while preserving all metadata.
        
        Args:
            input_path: Path to input JSON file
            output_path: Path to output JSON file
            
        Returns:
            Dictionary with translation statistics
        """
        print(f"Loading presentation from {input_path}...")
        with open(input_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        print(f"Translating to {self.target_language}...")
        if self.is_rtl:
            print(f"RTL Mode: ENABLED (Right-to-Left text direction will be applied)")
        print(f"Total slides: {data['total_slides']}")
        if 'slide_masters' in data:
            print(f"Slide masters: {len(data['slide_masters'])}")
        print("=" * 80)
        
        translated_data = {
            "presentation_name": data["presentation_name"],
            "total_slides": data["total_slides"],
            "target_language": self.target_language,
            "is_rtl": self.is_rtl,
            "slides": []
        }
        
        if "slide_masters" in data:
            translated_data["slide_masters"] = deepcopy(data["slide_masters"])
        
        start_time = time.time()
        for idx, slide in enumerate(data["slides"], 1):
            translated_slide = self.translate_slide(slide, idx)
            translated_data["slides"].append(translated_slide)
        
        elapsed_time = time.time() - start_time
        
        print("=" * 80)
        print(f"Saving translated presentation to {output_path}...")
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(translated_data, f, indent=2, ensure_ascii=False)
        
        print("\n" + "=" * 80)
        print("TRANSLATION COMPLETE!")
        print("=" * 80)
        print(f"Target language: {self.target_language}")
        if self.is_rtl:
            print(f"RTL Mode: ENABLED")
        print(f"Total slides translated: {data['total_slides']}")
        print(f"Total texts translated: {self.stats['total_texts_translated']}")
        print(f"API calls made: {self.stats['api_calls']}")
        print(f"Total tokens used: {self.stats['total_tokens_used']}")
        print(f"Time elapsed: {elapsed_time:.2f} seconds")
        print(f"Output saved to: {output_path}")
        print("=" * 80)
        
        return self.stats


# ============================================================================
# PART 3: PPTX REASSEMBLER
# ============================================================================

class PPTXReassembler:
    """
    Reassembles PowerPoint presentation from translated JSON.
    
    Strategy: Template-based approach
    - Uses original PowerPoint as template
    - Matches slides by layout information
    - Replaces only text content
    - Preserves all visual formatting automatically
    """
    
    def __init__(self, original_pptx_path: str, translated_json_path: str):
        """
        Initialize the reassembler.
        
        Args:
            original_pptx_path: Path to original PowerPoint file (template)
            translated_json_path: Path to translated JSON file
        """
        self.original_pptx_path = original_pptx_path
        self.translated_json_path = translated_json_path
        
        print(f"Loading original presentation: {original_pptx_path}")
        self.presentation = Presentation(original_pptx_path)
        
        print(f"Loading translated content: {translated_json_path}")
        with open(translated_json_path, 'r', encoding='utf-8') as f:
            self.translated_data = json.load(f)
        
        self.target_language = self.translated_data.get("target_language", "Unknown")
        self.is_rtl = self.translated_data.get("is_rtl", False)
        
        self.stats = {
            "slides_processed": 0,
            "elements_updated": 0,
            "text_runs_updated": 0,
            "tables_updated": 0,
            "charts_updated": 0,
            "notes_updated": 0,
            "rtl_paragraphs_set": 0,
            "shapes_mirrored": 0,
            "auto_shrink_enabled": 0
        }
        
        print(f"✓ Loaded {len(self.presentation.slides)} slides from original")
        print(f"✓ Loaded {len(self.translated_data['slides'])} slides from JSON")
    
    def find_shape_by_id(self, slide, shape_id: int):
        """
        Find a shape in a slide by its shape_id.
        
        Args:
            slide: PowerPoint slide object
            shape_id: Shape ID to find
            
        Returns:
            Shape object or None if not found
        """
        for shape in slide.shapes:
            if shape.shape_id == shape_id:
                return shape
        return None
    
    def set_paragraph_rtl(self, paragraph):
        """
        Set paragraph to Right-to-Left text direction.
        
        Args:
            paragraph: PowerPoint paragraph object
        """
        try:
            pPr = paragraph._element.get_or_add_pPr()
            pPr.set('rtl', '1')
            self.stats["rtl_paragraphs_set"] += 1
        except Exception as e:
            pass
    
    def mirror_shape_position(self, shape, slide_width):
        """
        Mirror shape position horizontally for RTL layout.
        
        Args:
            shape: PowerPoint shape object
            slide_width: Width of the slide
        """
        try:
            original_left = shape.left
            shape.left = slide_width - original_left - shape.width
            self.stats["shapes_mirrored"] += 1
        except Exception as e:
            pass
    
    def mirror_slide_layout(self, slide):
        """
        Mirror all shapes in a slide for RTL languages.
        
        Args:
            slide: PowerPoint slide object
        """
        if not self.is_rtl:
            return
        
        try:
            slide_width = self.presentation.slide_width
            
            for shape in slide.shapes:
                self.mirror_shape_position(shape, slide_width)
        except Exception as e:
            pass
    
    def enable_auto_shrink(self, text_frame):
        """
        Enable auto-shrink on text frame to prevent overflow.
        
        Args:
            text_frame: PowerPoint text frame object
        """
        try:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            self.stats["auto_shrink_enabled"] += 1
        except:
            try:
                text_frame.word_wrap = True
                text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            except:
                pass
    
    def update_text_runs(self, paragraph, translated_runs: list):
        """
        Update text runs in a paragraph with translated text.
        Preserves all formatting.
        
        Args:
            paragraph: PowerPoint paragraph object
            translated_runs: List of translated run dictionaries
        """
        existing_runs = list(paragraph.runs)
        
        if len(existing_runs) == len(translated_runs):
            for idx, (ppt_run, json_run) in enumerate(zip(existing_runs, translated_runs)):
                ppt_run.text = json_run.get("text", "")
                self.stats["text_runs_updated"] += 1
        else:
            while len(paragraph.runs) > 1:
                paragraph._element.remove(paragraph.runs[-1]._element)
            
            if len(paragraph.runs) == 0:
                paragraph.add_run()
            
            if len(translated_runs) > 0:
                paragraph.runs[0].text = translated_runs[0].get("text", "")
                self.stats["text_runs_updated"] += 1
            
            for json_run in translated_runs[1:]:
                new_run = paragraph.add_run()
                new_run.text = json_run.get("text", "")
                
                try:
                    font = new_run.font
                    if json_run.get("bold") is not None:
                        font.bold = json_run["bold"]
                    if json_run.get("italic") is not None:
                        font.italic = json_run["italic"]
                    if json_run.get("font_size"):
                        font.size = Pt(json_run["font_size"])
                except:
                    pass
                
                self.stats["text_runs_updated"] += 1
    
    def update_text_frame(self, shape, translated_element: dict):
        """
        Update text in a text frame (TextBox, AutoShape, etc.).
        
        Args:
            shape: PowerPoint shape with text_frame
            translated_element: Translated element dictionary from JSON
        """
        if not shape.has_text_frame:
            return
        
        text_frame = shape.text_frame
        translated_paragraphs = translated_element.get("paragraphs", [])
        
        existing_paragraphs = list(text_frame.paragraphs)
        
        for idx, translated_para in enumerate(translated_paragraphs):
            if idx < len(existing_paragraphs):
                ppt_para = existing_paragraphs[idx]
                translated_runs = translated_para.get("runs", [])
                self.update_text_runs(ppt_para, translated_runs)
                
                if self.is_rtl:
                    self.set_paragraph_rtl(ppt_para)
            else:
                ppt_para = text_frame.add_paragraph()
                translated_runs = translated_para.get("runs", [])
                self.update_text_runs(ppt_para, translated_runs)
                
                if self.is_rtl:
                    self.set_paragraph_rtl(ppt_para)
        
        while len(text_frame.paragraphs) > len(translated_paragraphs):
            try:
                text_frame._element.remove(text_frame.paragraphs[-1]._element)
            except:
                break
        
        self.enable_auto_shrink(text_frame)
    
    def update_table(self, shape, translated_table_data: dict):
        """
        Update table cell text with translations.
        
        Args:
            shape: PowerPoint shape with table
            translated_table_data: Translated table_data dictionary from JSON
        """
        if not shape.has_table:
            return
        
        table = shape.table
        translated_cells = translated_table_data.get("cells", [])
        
        for cell_data in translated_cells:
            row = cell_data.get("row")
            col = cell_data.get("column")
            
            if row is not None and col is not None:
                try:
                    cell = table.cell(row, col)
                    translated_paragraphs = cell_data.get("paragraphs", [])
                    
                    if cell.text_frame and translated_paragraphs:
                        for idx, translated_para in enumerate(translated_paragraphs):
                            if idx < len(cell.text_frame.paragraphs):
                                ppt_para = cell.text_frame.paragraphs[idx]
                                translated_runs = translated_para.get("runs", [])
                                self.update_text_runs(ppt_para, translated_runs)
                                
                                if self.is_rtl:
                                    self.set_paragraph_rtl(ppt_para)
                        
                        self.enable_auto_shrink(cell.text_frame)
                
                except Exception as e:
                    print(f"  Warning: Could not update cell ({row}, {col}): {e}")
        
        self.stats["tables_updated"] += 1
    
    def update_chart(self, shape, translated_chart_data: dict):
        """
        Update chart text elements (title, series names, categories).
        Note: Chart data values are NOT translated, only text labels.
        
        Args:
            shape: PowerPoint shape with chart
            translated_chart_data: Translated chart_data dictionary from JSON
        """
        if not shape.has_chart:
            return
        
        try:
            chart = shape.chart
            
            if translated_chart_data.get("title") and chart.has_title:
                try:
                    chart.chart_title.text_frame.text = translated_chart_data["title"]
                except Exception as e:
                    print(f"  Warning: Could not update chart title: {e}")
            
            translated_series = translated_chart_data.get("series", [])
            try:
                for idx, (ppt_series, json_series) in enumerate(zip(chart.series, translated_series)):
                    if json_series.get("name"):
                        ppt_series.name = json_series["name"]
            except Exception as e:
                print(f"  Warning: Could not update chart series names: {e}")
            
            self.stats["charts_updated"] += 1
            
        except Exception as e:
            print(f"  Warning: Error updating chart: {e}")
    
    def update_speaker_notes(self, slide, translated_notes: dict):
        """
        Update speaker notes with translated text.
        
        Args:
            slide: PowerPoint slide object
            translated_notes: Translated notes dictionary from JSON
        """
        try:
            if not slide.has_notes_slide:
                slide.notes_slide
            
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            
            translated_text = translated_notes.get("text", "")
            notes_text_frame.text = translated_text
            
            if self.is_rtl:
                for paragraph in notes_text_frame.paragraphs:
                    self.set_paragraph_rtl(paragraph)
            
            self.enable_auto_shrink(notes_text_frame)
            
            self.stats["notes_updated"] += 1
            
        except Exception as e:
            print(f"  Warning: Could not update speaker notes: {e}")
    
    def update_slide(self, slide_idx: int):
        """
        Update a single slide with translated content.
        
        Args:
            slide_idx: 0-based slide index
        """
        ppt_slide = self.presentation.slides[slide_idx]
        
        translated_slide = self.translated_data["slides"][slide_idx]
        
        print(f"Processing slide {slide_idx + 1}/{len(self.presentation.slides)}...", end=" ")
        
        self.mirror_slide_layout(ppt_slide)
        
        translated_elements = translated_slide.get("elements", [])
        
        for element in translated_elements:
            shape_id = element.get("shape_id")
            element_type = element.get("element_type")
            
            shape = self.find_shape_by_id(ppt_slide, shape_id)
            
            if not shape:
                continue
            
            try:
                if element_type in ["TextBox", "AutoShape"]:
                    self.update_text_frame(shape, element)
                    self.stats["elements_updated"] += 1
                
                elif element_type == "Table":
                    table_data = element.get("table_data")
                    if table_data:
                        self.update_table(shape, table_data)
                        self.stats["elements_updated"] += 1
                
                elif element_type == "Chart":
                    chart_data = element.get("chart_data")
                    if chart_data:
                        self.update_chart(shape, chart_data)
                        self.stats["elements_updated"] += 1
                
            except Exception as e:
                print(f"\n  Warning: Error updating shape {shape_id} ({element_type}): {e}")
        
        translated_notes = translated_slide.get("speaker_notes")
        if translated_notes:
            self.update_speaker_notes(ppt_slide, translated_notes)
        
        self.stats["slides_processed"] += 1
        print("✓")
    
    def verify_slide_count(self):
        """
        Verify that original PPT and translated JSON have same number of slides.
        
        Returns:
            bool: True if counts match
        """
        ppt_count = len(self.presentation.slides)
        json_count = len(self.translated_data["slides"])
        
        if ppt_count != json_count:
            print(f"\n⚠️  WARNING: Slide count mismatch!")
            print(f"   Original PPT: {ppt_count} slides")
            print(f"   Translated JSON: {json_count} slides")
            print(f"   Will process minimum: {min(ppt_count, json_count)} slides")
            return False
        
        return True
    
    def reassemble(self, output_path: str):
        """
        Main reassembly process.
        
        Args:
            output_path: Path to save the reassembled PowerPoint
        """
        print("\n" + "=" * 80)
        print("STARTING REASSEMBLY")
        print("=" * 80)
        print(f"Target language: {self.target_language}")
        if self.is_rtl:
            print(f"RTL mode: ENABLED (Right-to-Left text direction)")
        print(f"Auto-shrink: ENABLED (Text will auto-fit to prevent overflow)")
        
        self.verify_slide_count()
        
        num_slides = min(len(self.presentation.slides), len(self.translated_data["slides"]))
        
        print(f"\nProcessing {num_slides} slides...")
        print("-" * 80)
        
        for slide_idx in range(num_slides):
            self.update_slide(slide_idx)
        
        print("-" * 80)
        print(f"\nSaving reassembled presentation to: {output_path}")
        self.presentation.save(output_path)
        
        print("\n" + "=" * 80)
        print("REASSEMBLY COMPLETE!")
        print("=" * 80)
        print(f"Slides processed: {self.stats['slides_processed']}")
        print(f"Elements updated: {self.stats['elements_updated']}")
        print(f"Text runs updated: {self.stats['text_runs_updated']}")
        print(f"Tables updated: {self.stats['tables_updated']}")
        print(f"Charts updated: {self.stats['charts_updated']}")
        print(f"Speaker notes updated: {self.stats['notes_updated']}")
        if self.is_rtl:
            print(f"RTL paragraphs set: {self.stats['rtl_paragraphs_set']}")
            print(f"Shapes mirrored (layout flipped): {self.stats['shapes_mirrored']}")
        print(f"Auto-shrink enabled: {self.stats['auto_shrink_enabled']} text frames")
        print(f"\n✓ Output saved to: {output_path}")
        print("=" * 80)
        
        return self.stats


# ============================================================================
# MAIN PIPELINE INTERFACE
# ============================================================================

def run_full_pipeline(pptx_file: str, target_language: str = "Spanish", 
                     api_key: str = None, output_dir: str = None):
    """
    Run the complete pipeline: Extract → Translate → Reassemble
    
    Args:
        pptx_file: Path to input PowerPoint file
        target_language: Target language for translation
        api_key: OpenAI API key (optional, reads from .env if not provided)
        output_dir: Output directory (optional, uses same dir as input if not provided)
    
    Returns:
        dict: Statistics from all three stages
    """
    base_name = os.path.splitext(os.path.basename(pptx_file))[0]
    
    if output_dir is None:
        output_dir = os.path.dirname(pptx_file) or "."
    
    extracted_json = os.path.join(output_dir, f"{base_name}_extracted.json")
    translated_json = os.path.join(output_dir, f"{base_name}_translated_{target_language.lower()}.json")
    output_pptx = os.path.join(output_dir, f"{base_name}_{target_language.lower()}.pptx")
    
    print("\n" + "=" * 80)
    print("POWERPOINT TRANSLATION PIPELINE")
    print("=" * 80)
    print(f"Input file: {pptx_file}")
    print(f"Target language: {target_language}")
    print(f"Output directory: {output_dir}")
    print("=" * 80 + "\n")
    
    # Stage 1: Extract
    print("STAGE 1: EXTRACTION")
    print("-" * 80)
    extractor = PPTXExtractor(pptx_file)
    extracted_data = extractor.extract_all()
    extractor.save_to_json(extracted_json)
    print(f"✓ Extraction complete: {len(extracted_data['slides'])} slides extracted")
    print()
    
    # Stage 2: Translate
    print("STAGE 2: TRANSLATION")
    print("-" * 80)
    translator = PPTTranslator(api_key=api_key, target_language=target_language)
    translation_stats = translator.translate_presentation(extracted_json, translated_json)
    print()
    
    # Stage 3: Reassemble
    print("STAGE 3: REASSEMBLY")
    print("-" * 80)
    reassembler = PPTXReassembler(pptx_file, translated_json)
    reassembly_stats = reassembler.reassemble(output_pptx)
    
    # Final summary
    print("\n" + "=" * 80)
    print("PIPELINE COMPLETE!")
    print("=" * 80)
    print(f"✓ Extracted: {extracted_json}")
    print(f"✓ Translated: {translated_json}")
    print(f"✓ Final output: {output_pptx}")
    print("=" * 80 + "\n")
    
    return {
        "extraction": {"slides": len(extracted_data['slides'])},
        "translation": translation_stats,
        "reassembly": reassembly_stats
    }


def main():
    """Main CLI interface"""
    parser = argparse.ArgumentParser(
        description="PowerPoint Translation Assembly Pipeline",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Run full pipeline (extract + translate + reassemble)
  python ppt-translation-assembly-pipeline.py presentation.pptx --language Spanish
  
  # Extract only
  python ppt-translation-assembly-pipeline.py presentation.pptx --extract-only
  
  # Translate only (requires extracted JSON)
  python ppt-translation-assembly-pipeline.py extracted.json --translate-only --language French
  
  # Reassemble only (requires original PPTX and translated JSON)
  python ppt-translation-assembly-pipeline.py original.pptx translated.json --reassemble-only
  
  # Custom output directory
  python ppt-translation-assembly-pipeline.py presentation.pptx --language Arabic --output ./outputs

Features:
  - Comprehensive metadata preservation
  - SmartArt support with hierarchical extraction
  - RTL language support (Arabic, Hebrew, etc.)
  - Automatic layout mirroring for RTL
  - Table, Chart, and complex shape handling
  - Speaker notes translation
        """
    )
    
    parser.add_argument("input_file", help="Input PowerPoint file or JSON file")
    parser.add_argument("input_file2", nargs='?', help="Second input file (for reassemble-only mode)")
    parser.add_argument("-l", "--language", default="Spanish", help="Target language (default: Spanish)")
    parser.add_argument("-k", "--api-key", help="OpenAI API key (default: from .env)")
    parser.add_argument("-o", "--output", help="Output directory (default: same as input)")
    parser.add_argument("--extract-only", action="store_true", help="Only extract content to JSON")
    parser.add_argument("--translate-only", action="store_true", help="Only translate JSON file")
    parser.add_argument("--reassemble-only", action="store_true", help="Only reassemble from JSON")
    
    args = parser.parse_args()
    
    try:
        if args.extract_only:
            # Extract only mode
            pptx_file = args.input_file
            base_name = os.path.splitext(os.path.basename(pptx_file))[0]
            output_dir = args.output or os.path.dirname(pptx_file) or "."
            output_json = os.path.join(output_dir, f"{base_name}_extracted.json")
            
            print("Running extraction only...")
            extractor = PPTXExtractor(pptx_file)
            extractor.extract_all()
            extractor.save_to_json(output_json)
            
        elif args.translate_only:
            # Translate only mode
            json_file = args.input_file
            base_name = os.path.splitext(os.path.basename(json_file))[0]
            output_dir = args.output or os.path.dirname(json_file) or "."
            output_json = os.path.join(output_dir, f"{base_name}_translated_{args.language.lower()}.json")
            
            print("Running translation only...")
            translator = PPTTranslator(api_key=args.api_key, target_language=args.language)
            translator.translate_presentation(json_file, output_json)
            
        elif args.reassemble_only:
            # Reassemble only mode
            if not args.input_file2:
                print("Error: Reassemble mode requires two arguments: original.pptx and translated.json")
                return 1
            
            pptx_file = args.input_file
            json_file = args.input_file2
            base_name = os.path.splitext(os.path.basename(json_file))[0]
            output_dir = args.output or os.path.dirname(pptx_file) or "."
            output_pptx = os.path.join(output_dir, f"{base_name}_reassembled.pptx")
            
            print("Running reassembly only...")
            reassembler = PPTXReassembler(pptx_file, json_file)
            reassembler.reassemble(output_pptx)
            
        else:
            # Full pipeline mode
            pptx_file = args.input_file
            run_full_pipeline(
                pptx_file=pptx_file,
                target_language=args.language,
                api_key=args.api_key,
                output_dir=args.output
            )
        
        return 0
        
    except Exception as e:
        print(f"\n❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    exit(main())